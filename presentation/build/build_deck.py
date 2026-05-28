"""
Build the Restylane clinical reference deck as a .pptx.

Style: AAA Academy-inspired medical infographic. Bilingual Hebrew/English.
All visuals are drawn as native PowerPoint shapes (no external images required)
so the deck is fully self-contained and edits cleanly in PowerPoint/Keynote.
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree

# ─────────────────────────────────────────────────────────────────────────────
# Bidi safety: keep digit / Latin runs in their natural LTR order inside RTL
# paragraphs. Use Unicode embedding marks so PowerPoint and LibreOffice both
# render numbers correctly.
# ─────────────────────────────────────────────────────────────────────────────

_LRM = "‎"  # LEFT-TO-RIGHT MARK (invisible)


def bidi(text: str) -> str:
    """In RTL paragraphs, insert invisible LTR marks around Latin/digit
    runs so neutral chars (hyphens, slashes) inherit the LTR direction
    and the cluster doesn't get re-ordered."""
    if not text:
        return text
    pattern = re.compile(r"[0-9A-Za-z][0-9A-Za-z\-–_+'/&.:%× ]*[0-9A-Za-z)']?")
    return pattern.sub(lambda m: f"{_LRM}{m.group(0)}{_LRM}", text)


_HEB_RE = re.compile(r"[֐-׿]")


def _is_hebrew_run(s: str) -> bool:
    return bool(_HEB_RE.search(s))


def _segment_bidi(text: str):
    """Yield (segment, is_hebrew) tuples splitting on script changes.
    Digits and Latin letters always start a new LTR segment so they get
    proper LTR ordering/metrics in mixed-script Hebrew paragraphs."""
    if not text:
        return
    cur = []
    cur_is_heb = None
    NEUTRAL = "-–—_+()/[]{}.,:;%×&!?\"'"
    for ch in text:
        if _HEB_RE.match(ch):
            is_heb = True
        elif ch.isdigit() or (ch.isascii() and ch.isalpha()):
            # Strong LTR: digits and Latin letters always force a new LTR run
            is_heb = False
        elif ch in NEUTRAL or ch == " ":
            # Neutral: inherit current direction. If first char, default
            # to RTL (Hebrew paragraph) so leading punctuation stays attached.
            is_heb = cur_is_heb if cur_is_heb is not None else True
        else:
            is_heb = cur_is_heb if cur_is_heb is not None else False
        if cur_is_heb is None:
            cur_is_heb = is_heb
        if is_heb != cur_is_heb:
            yield "".join(cur), cur_is_heb
            cur = [ch]
            cur_is_heb = is_heb
        else:
            cur.append(ch)
    if cur:
        yield "".join(cur), cur_is_heb

# ─────────────────────────────────────────────────────────────────────────────
# Brand palette (Restylane / AAA Academy inspired)
# ─────────────────────────────────────────────────────────────────────────────

CREAM       = RGBColor(0xF5, 0xED, 0xE3)
CREAM_DEEP  = RGBColor(0xE8, 0xD9, 0xC5)
SAND        = RGBColor(0xD9, 0xC2, 0xA3)
TAUPE       = RGBColor(0xA8, 0x8D, 0x6F)
GOLD        = RGBColor(0xB8, 0x95, 0x6A)
GOLD_DARK   = RGBColor(0x8A, 0x6E, 0x4A)
INK         = RGBColor(0x2C, 0x1F, 0x15)
INK_SOFT    = RGBColor(0x4A, 0x3A, 0x2A)
LINE_SOFT   = RGBColor(0xC9, 0xB8, 0xA0)
WHITE_WARM  = RGBColor(0xFA, 0xF6, 0xF0)
WHITE_PURE  = RGBColor(0xFF, 0xFF, 0xFF)

NASHA_BLUE     = RGBColor(0x2E, 0x4F, 0x8A)
NASHA_BLUE_LT  = RGBColor(0x6E, 0x8F, 0xC4)
OBT_PURPLE     = RGBColor(0x7A, 0x5A, 0x99)
OBT_PURPLE_LT  = RGBColor(0xC3, 0xAE, 0xD6)
ACCENT_ROSE    = RGBColor(0xC8, 0x6D, 0x6D)
ACCENT_TEAL    = RGBColor(0x4E, 0x96, 0x96)

DANGER         = RGBColor(0xB2, 0x3B, 0x3B)
SUCCESS        = RGBColor(0x3F, 0x8C, 0x55)
WARN           = RGBColor(0xC4, 0x84, 0x2C)

FONT_HEB = "Noto Sans Hebrew"
FONT_LAT = "Noto Sans"
FONT_SERIF = "Noto Serif Hebrew"

# Slide geometry (16:9 widescreen, 13.333 x 7.5 inches)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─────────────────────────────────────────────────────────────────────────────
# Low-level XML helpers (gradient fills, RTL paragraph direction, etc.)
# ─────────────────────────────────────────────────────────────────────────────

def _solid_fill(elem, color: RGBColor):
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    fill = etree.SubElement(elem, qn("a:solidFill"))
    clr = etree.SubElement(fill, qn("a:srgbClr"))
    clr.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
    return fill


def set_slide_bg_gradient(slide, top: RGBColor, bottom: RGBColor, angle: int = 5400000):
    """Fill the slide background with a two-stop linear gradient."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()  # creates default gradient
    stops = fill.gradient_stops
    stops[0].color.rgb = top
    stops[1].color.rgb = bottom
    try:
        fill.gradient_angle = angle / 60000.0  # python-pptx wants degrees
    except Exception:
        pass


def set_slide_bg_solid(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_shape_gradient(shape, top: RGBColor, bottom: RGBColor, angle: float = 90.0):
    """Force a two-stop gradient onto an autoshape, replacing any fill."""
    spPr = shape.fill._xPr.find(qn("p:spPr")) or shape.line._get_or_add_ln().getparent().find(qn("p:spPr"))
    if spPr is None:
        spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    # remove existing fill elements
    for tag in ("a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:noFill"):
        for old in spPr.findall(qn(tag)):
            spPr.remove(old)
    grad = etree.SubElement(spPr, qn("a:gradFill"))
    grad.set("rotWithShape", "1")
    lst = etree.SubElement(grad, qn("a:gsLst"))
    for pos, color in ((0, top), (100000, bottom)):
        gs = etree.SubElement(lst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
    lin = etree.SubElement(grad, qn("a:lin"))
    lin.set("ang", str(int(angle * 60000)))
    lin.set("scaled", "0")
    # ensure spPr is in the right place
    parent = spPr.getparent()
    if parent is not None and spPr.getnext() is not None:
        pass


def set_paragraph_rtl(paragraph, rtl: bool = True):
    """Mark a paragraph as right-to-left for proper Hebrew rendering."""
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1" if rtl else "0")


def add_textbox(
    slide,
    left, top, width, height,
    text: str = "",
    font: str = FONT_HEB,
    size: int = 18,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = INK,
    align: str = "right",
    rtl: bool = True,
    anchor: str = "top",
    line_spacing: float = 1.15,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.TOP)

    lines = text.split("\n") if text else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {
            "right": PP_ALIGN.RIGHT,
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "justify": PP_ALIGN.JUSTIFY,
        }.get(align, PP_ALIGN.RIGHT)
        if line_spacing:
            p.line_spacing = line_spacing
        # If line has no Hebrew, treat it as LTR even when the box default
        # is RTL — avoids digit-only strings ("12–18") getting reversed.
        effective_rtl = rtl and _is_hebrew_run(line)
        set_paragraph_rtl(p, effective_rtl)

        # Split this line into Hebrew vs non-Hebrew segments so each can use
        # its own font and the bidi engine sees clean direction transitions.
        segments = list(_segment_bidi(line)) if effective_rtl and line else [(line, False)]
        if not segments:
            segments = [("", False)]
        for seg_text, seg_is_heb in segments:
            run = p.add_run()
            run.text = seg_text
            seg_font = font if (font == FONT_HEB or seg_is_heb) else font
            # Use Hebrew font for Hebrew runs, Latin font for Latin/digit runs
            if seg_is_heb:
                seg_font = FONT_HEB if font in (FONT_HEB, FONT_SERIF) else font
            else:
                seg_font = FONT_LAT if font == FONT_HEB else font
            run.font.name = seg_font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            # Set the appropriate font slot in rPr
            rPr = run._r.get_or_add_rPr()
            if seg_is_heb:
                # Mark as complex script so the engine uses the cs typeface
                cs = rPr.find(qn("a:cs"))
                if cs is None:
                    cs = etree.SubElement(rPr, qn("a:cs"))
                cs.set("typeface", FONT_HEB)
                # Also set the regular typeface (PowerPoint fallback)
                latin = rPr.find(qn("a:latin"))
                if latin is None:
                    latin = etree.SubElement(rPr, qn("a:latin"))
                latin.set("typeface", FONT_HEB)
            else:
                # Force regular Latin typeface AND complex-script to the same
                # font so LibreOffice can't substitute a different cs font
                # (which causes weird digit sizes inside Hebrew paragraphs).
                latin = rPr.find(qn("a:latin"))
                if latin is None:
                    latin = etree.SubElement(rPr, qn("a:latin"))
                latin.set("typeface", seg_font)
                cs = rPr.find(qn("a:cs"))
                if cs is None:
                    cs = etree.SubElement(rPr, qn("a:cs"))
                cs.set("typeface", seg_font)
                ea = rPr.find(qn("a:ea"))
                if ea is None:
                    ea = etree.SubElement(rPr, qn("a:ea"))
                ea.set("typeface", seg_font)
    return tb


def add_rect(slide, left, top, width, height, fill: RGBColor | None = None,
             line: RGBColor | None = None, line_w: float = 0.5,
             shape_type=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    return s


def add_line(slide, x1, y1, x2, y2, color: RGBColor = LINE_SOFT, weight: float = 1.0):
    s = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    s.line.color.rgb = color
    s.line.width = Pt(weight)
    return s


def add_oval(slide, left, top, width, height, fill=None, line=None, line_w=0.75):
    return add_rect(slide, left, top, width, height, fill=fill, line=line,
                    line_w=line_w, shape_type=MSO_SHAPE.OVAL)


def add_triangle(slide, left, top, width, height, fill=None, line=None, line_w=2.0):
    return add_rect(slide, left, top, width, height, fill=fill, line=line,
                    line_w=line_w, shape_type=MSO_SHAPE.ISOSCELES_TRIANGLE)


# ─────────────────────────────────────────────────────────────────────────────
# Composed visuals
# ─────────────────────────────────────────────────────────────────────────────

def draw_product_visual(slide, panel_l, panel_t, panel_w, panel_h,
                        kind: str, color: RGBColor, color_lt: RGBColor):
    """Draw a clean abstract icon for a product, centered in the given panel."""
    cx = panel_l + panel_w // 2
    cy = panel_t + panel_h // 2 - Inches(0.2)

    # Inner subtle circle background (gives the icon a 'medallion' feel)
    medal_w = Inches(3.2)
    medal_h = Inches(3.2)
    add_oval(slide, cx - medal_w // 2, cy - medal_h // 2, medal_w, medal_h,
             fill=CREAM, line=color_lt, line_w=1.0)
    # second ring
    ring_w = Inches(3.6)
    ring_h = Inches(3.6)
    add_oval(slide, cx - ring_w // 2, cy - ring_h // 2, ring_w, ring_h,
             fill=None, line=color_lt, line_w=0.5)

    if kind == "triangle_jaw":   # Lyft → Bone Mimicry
        # Large outlined triangle suggesting jawline geometry
        tri_w = Inches(2.2)
        tri_h = Inches(2.0)
        # Custom triangle using freeform - use Isoceles_triangle, then rotate
        t = add_triangle(slide, cx - tri_w // 2, cy - tri_h // 2,
                          tri_w, tri_h, fill=None, line=color, line_w=4.0)
        # Rotate -25° for a jawline feel
        t.rotation = -25
        # Inner thinner triangle echo
        inner_w = Inches(1.4)
        inner_h = Inches(1.3)
        t2 = add_triangle(slide, cx - inner_w // 2, cy - inner_h // 2 + Inches(0.05),
                          inner_w, inner_h, fill=None, line=color_lt, line_w=2.0)
        t2.rotation = -25
        return "Bone Mimicry"
    elif kind == "mesh_jaw":     # Defyne → Tissue Mimicry
        # Curved mesh pattern (grid)
        mw = Inches(2.6)
        mh = Inches(2.0)
        ml = cx - mw // 2
        mt = cy - mh // 2 + Inches(0.1)
        # Rounded background patch
        add_rect(slide, ml - Inches(0.2), mt - Inches(0.2),
                 mw + Inches(0.4), mh + Inches(0.4),
                 fill=color_lt, line=None,
                 shape_type=MSO_SHAPE.OVAL)
        # mesh lines
        rows, cols = 8, 12
        for r in range(rows + 1):
            add_line(slide, ml, mt + Emu(int(mh.emu * r / rows)),
                     ml + mw, mt + Emu(int(mh.emu * r / rows)),
                     color=color, weight=0.75)
        for c in range(cols + 1):
            add_line(slide, ml + Emu(int(mw.emu * c / cols)), mt,
                     ml + Emu(int(mw.emu * c / cols)), mt + mh,
                     color=color, weight=0.75)
        return "Tissue Mimicry"
    elif kind == "soft_cheek":   # Volyme → soft volume
        # Layered glowing circles in lighter→darker
        sizes = [(2.4, color_lt), (1.7, color_lt), (1.0, color)]
        for sz, c in sizes:
            sw = Inches(sz)
            add_oval(slide, cx - sw // 2, cy - sw // 2, sw, sw,
                     fill=c, line=None)
        return "Soft Volume"
    elif kind == "lips":         # Kysse
        # Stylized lips: two overlapping ovals (upper + lower lip)
        lw = Inches(2.6)
        lh_top = Inches(0.55)
        lh_bot = Inches(0.7)
        # Upper lip with cupid's bow hint
        add_oval(slide, cx - lw // 2, cy - Inches(0.45), lw, lh_top,
                 fill=color, line=None)
        # Lower lip (fuller)
        add_oval(slide, cx - lw // 2, cy + Inches(0.05), lw, lh_bot,
                 fill=color_lt, line=color, line_w=1.0)
        # Cupid bow notch
        notch_w = Inches(0.45)
        notch_h = Inches(0.25)
        add_oval(slide, cx - notch_w // 2, cy - Inches(0.55),
                 notch_w, notch_h, fill=CREAM, line=None)
        return "Lip Dynamics"
    elif kind == "undereye":     # Eyelight
        # Eye almond + crescent under it
        eye_w = Inches(1.6)
        eye_h = Inches(0.7)
        # Eye outline
        add_oval(slide, cx - eye_w // 2, cy - Inches(0.5),
                 eye_w, eye_h, fill=WHITE_WARM, line=color, line_w=1.5)
        # Iris
        iris_w = Inches(0.55)
        add_oval(slide, cx - iris_w // 2, cy - Inches(0.4),
                 iris_w, iris_w, fill=color, line=None)
        # Crescent under-eye (Tear-trough zone)
        cresc_w = Inches(2.0)
        cresc_h = Inches(0.6)
        add_oval(slide, cx - cresc_w // 2, cy + Inches(0.4),
                 cresc_w, cresc_h, fill=color_lt, line=color, line_w=1.0)
        # Shadow line just below eye
        add_line(slide, cx - Inches(0.85), cy + Inches(0.35),
                 cx + Inches(0.85), cy + Inches(0.35),
                 color=color, weight=1.0)
        return "Tear-Trough"
    elif kind == "wrinkle_line":  # Refyne / Classic
        # Three curved 'wrinkle' lines + a smooth highlight under them
        for i in range(3):
            y = cy - Inches(0.45) + Inches(0.45) * i
            add_line(slide, cx - Inches(1.3), y,
                     cx + Inches(1.3), y + Inches(0.15) - Inches(0.05) * i,
                     color=color, weight=2.0)
        # subtle smooth line underneath (effect after treatment)
        add_line(slide, cx - Inches(1.3), cy + Inches(0.9),
                 cx + Inches(1.3), cy + Inches(0.9),
                 color=color_lt, weight=4.0)
        return "Dynamic Support"
    elif kind == "dots":          # Vital / Vital Light
        import random
        random.seed(11)
        # Concentric ring of micro-dots
        for r_in in (0.6, 1.0, 1.4):
            n = int(r_in * 12)
            from math import pi, cos, sin
            for i in range(n):
                t = 2 * pi * i / n
                x = cx + Emu(int(Inches(r_in).emu * cos(t)))
                y = cy + Emu(int(Inches(r_in).emu * sin(t)))
                add_oval(slide, x - Inches(0.05), y - Inches(0.05),
                         Inches(0.1), Inches(0.1), fill=color, line=None)
        # Center
        add_oval(slide, cx - Inches(0.15), cy - Inches(0.15),
                 Inches(0.3), Inches(0.3), fill=color, line=None)
        return "Biorevitalization"
    return ""


def draw_face_silhouette_profile(slide, cx, cy, scale=1.0, facing="left",
                                 skin_top=RGBColor(0xF1, 0xDC, 0xC4),
                                 skin_bot=RGBColor(0xE0, 0xC0, 0xA0)):
    """
    Draw a stylized side-profile silhouette using grouped shapes.
    cx, cy = approx center of the face in EMU.
    scale = relative size (1.0 ≈ 4.0in tall)
    facing = 'left' or 'right' (where the nose points)
    Returns the list of shapes drawn.
    """
    shapes = []
    sign = -1 if facing == "left" else 1

    # Hair / head back (large rounded shape)
    head_w = Emu(int(Inches(2.4).emu * scale))
    head_h = Emu(int(Inches(3.2).emu * scale))
    head_l = cx - head_w // 2 - Emu(int(sign * Inches(0.15).emu * scale))
    head_t = cy - head_h // 2 - Emu(int(Inches(0.2).emu * scale))
    head = add_rect(slide, head_l, head_t, head_w, head_h,
                    fill=RGBColor(0x6B, 0x4F, 0x3A), shape_type=MSO_SHAPE.OVAL)
    shapes.append(head)

    # Face oval (lighter, offset toward facing direction)
    face_w = Emu(int(Inches(1.7).emu * scale))
    face_h = Emu(int(Inches(2.4).emu * scale))
    face_l = cx - face_w // 2 + Emu(int(sign * Inches(-0.35).emu * scale))
    face_t = cy - face_h // 2 + Emu(int(Inches(0.05).emu * scale))
    face = add_rect(slide, face_l, face_t, face_w, face_h,
                    fill=skin_top, shape_type=MSO_SHAPE.OVAL)
    shapes.append(face)

    # Cheek shading
    cheek_w = Emu(int(Inches(0.9).emu * scale))
    cheek_h = Emu(int(Inches(0.8).emu * scale))
    cheek_l = cx - cheek_w // 2 + Emu(int(sign * Inches(-0.2).emu * scale))
    cheek_t = cy - cheek_h // 2 + Emu(int(Inches(0.1).emu * scale))
    cheek = add_rect(slide, cheek_l, cheek_t, cheek_w, cheek_h,
                     fill=skin_bot, shape_type=MSO_SHAPE.OVAL)
    cheek.fill.transparency = 0.5  # may be ignored by viewer
    shapes.append(cheek)

    # Neck
    neck_w = Emu(int(Inches(0.8).emu * scale))
    neck_h = Emu(int(Inches(1.2).emu * scale))
    neck_l = cx - neck_w // 2 + Emu(int(sign * Inches(-0.1).emu * scale))
    neck_t = cy + face_h // 2 - Emu(int(Inches(0.2).emu * scale))
    neck = add_rect(slide, neck_l, neck_t, neck_w, neck_h,
                    fill=skin_top, shape_type=MSO_SHAPE.OVAL)
    shapes.append(neck)

    # Shoulder
    sh_w = Emu(int(Inches(3.0).emu * scale))
    sh_h = Emu(int(Inches(1.2).emu * scale))
    sh_l = cx - sh_w // 2 + Emu(int(sign * Inches(-0.2).emu * scale))
    sh_t = cy + face_h // 2 + Emu(int(Inches(0.7).emu * scale))
    sh = add_rect(slide, sh_l, sh_t, sh_w, sh_h,
                  fill=INK_SOFT, shape_type=MSO_SHAPE.OVAL)
    shapes.append(sh)

    # Approximate landmarks for overlays (in EMU)
    landmarks = {
        "ear":       (cx + Emu(int(sign * -Inches(0.1).emu * scale)),
                      cy + Emu(int(Inches(0.0).emu * scale))),
        "jaw_angle": (cx + Emu(int(sign * -Inches(0.05).emu * scale)),
                      cy + Emu(int(Inches(0.85).emu * scale))),
        "chin":      (cx + Emu(int(sign * -Inches(0.9).emu * scale)),
                      cy + Emu(int(Inches(1.1).emu * scale))),
        "cheek":     (cx + Emu(int(sign * -Inches(0.55).emu * scale)),
                      cy + Emu(int(Inches(0.2).emu * scale))),
        "temple":    (cx + Emu(int(sign * -Inches(0.45).emu * scale)),
                      cy - Emu(int(Inches(0.9).emu * scale))),
        "lip":       (cx + Emu(int(sign * -Inches(0.95).emu * scale)),
                      cy + Emu(int(Inches(0.55).emu * scale))),
        "undereye":  (cx + Emu(int(sign * -Inches(0.7).emu * scale)),
                      cy - Emu(int(Inches(0.3).emu * scale))),
        "nose":      (cx + Emu(int(sign * -Inches(1.05).emu * scale)),
                      cy + Emu(int(Inches(0.05).emu * scale))),
    }
    return shapes, landmarks


def draw_face_silhouette_front(slide, cx, cy, scale=1.0):
    """Frontal face silhouette for the face-map slide."""
    head_w = Emu(int(Inches(3.0).emu * scale))
    head_h = Emu(int(Inches(3.8).emu * scale))
    head_l = cx - head_w // 2
    head_t = cy - head_h // 2
    head = add_rect(slide, head_l, head_t, head_w, head_h,
                    fill=RGBColor(0xF1, 0xDC, 0xC4), line=TAUPE, line_w=1.0,
                    shape_type=MSO_SHAPE.OVAL)

    # Hair cap
    hair_w = Emu(int(Inches(3.2).emu * scale))
    hair_h = Emu(int(Inches(1.4).emu * scale))
    hair_l = cx - hair_w // 2
    hair_t = cy - head_h // 2 - Emu(int(Inches(0.2).emu * scale))
    add_rect(slide, hair_l, hair_t, hair_w, hair_h,
             fill=RGBColor(0x6B, 0x4F, 0x3A), shape_type=MSO_SHAPE.OVAL)

    # Re-draw face oval on top of hair for clean forehead
    face = add_rect(slide, head_l, head_t + Emu(int(Inches(0.55).emu*scale)),
                    head_w, head_h - Emu(int(Inches(0.4).emu*scale)),
                    fill=RGBColor(0xF1, 0xDC, 0xC4), line=None,
                    shape_type=MSO_SHAPE.OVAL)

    # Brows
    brow_w = Emu(int(Inches(0.5).emu * scale))
    brow_h = Emu(int(Inches(0.08).emu * scale))
    for side in (-1, 1):
        bl = cx + Emu(int(side * Inches(0.55).emu * scale)) - brow_w // 2
        bt = cy - Emu(int(Inches(0.5).emu * scale))
        add_rect(slide, bl, bt, brow_w, brow_h, fill=INK_SOFT,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

    # Eyes (small ovals)
    eye_w = Emu(int(Inches(0.32).emu * scale))
    eye_h = Emu(int(Inches(0.18).emu * scale))
    for side in (-1, 1):
        el = cx + Emu(int(side * Inches(0.55).emu * scale)) - eye_w // 2
        et = cy - Emu(int(Inches(0.3).emu * scale))
        add_rect(slide, el, et, eye_w, eye_h, fill=WHITE_PURE, line=INK_SOFT,
                 line_w=0.5, shape_type=MSO_SHAPE.OVAL)
        # pupil
        pw = Emu(int(Inches(0.12).emu * scale))
        ph = Emu(int(Inches(0.12).emu * scale))
        add_rect(slide, el + (eye_w - pw)//2, et + (eye_h - ph)//2, pw, ph,
                 fill=INK, shape_type=MSO_SHAPE.OVAL)

    # Nose hint (thin vertical line)
    add_line(slide, cx, cy - Emu(int(Inches(0.1).emu * scale)),
             cx, cy + Emu(int(Inches(0.35).emu * scale)),
             color=TAUPE, weight=1.0)

    # Lips
    lip_w = Emu(int(Inches(0.6).emu * scale))
    lip_h = Emu(int(Inches(0.12).emu * scale))
    add_rect(slide, cx - lip_w//2, cy + Emu(int(Inches(0.6).emu * scale)),
             lip_w, lip_h, fill=RGBColor(0xC0, 0x80, 0x80),
             shape_type=MSO_SHAPE.OVAL)

    # Neck
    nw = Emu(int(Inches(1.0).emu * scale))
    nh = Emu(int(Inches(0.8).emu * scale))
    add_rect(slide, cx - nw//2, cy + Emu(int(Inches(1.7).emu * scale)), nw, nh,
             fill=RGBColor(0xF1, 0xDC, 0xC4), shape_type=MSO_SHAPE.OVAL)

    return {
        "tear_trough_L": (cx - Emu(int(Inches(0.55).emu*scale)), cy - Emu(int(Inches(0.1).emu*scale))),
        "tear_trough_R": (cx + Emu(int(Inches(0.55).emu*scale)), cy - Emu(int(Inches(0.1).emu*scale))),
        "cheek_L":       (cx - Emu(int(Inches(1.05).emu*scale)), cy + Emu(int(Inches(0.05).emu*scale))),
        "cheek_R":       (cx + Emu(int(Inches(1.05).emu*scale)), cy + Emu(int(Inches(0.05).emu*scale))),
        "temple_L":      (cx - Emu(int(Inches(1.25).emu*scale)), cy - Emu(int(Inches(0.7).emu*scale))),
        "temple_R":      (cx + Emu(int(Inches(1.25).emu*scale)), cy - Emu(int(Inches(0.7).emu*scale))),
        "nlf_L":         (cx - Emu(int(Inches(0.55).emu*scale)), cy + Emu(int(Inches(0.4).emu*scale))),
        "nlf_R":         (cx + Emu(int(Inches(0.55).emu*scale)), cy + Emu(int(Inches(0.4).emu*scale))),
        "lips":          (cx, cy + Emu(int(Inches(0.65).emu*scale))),
        "marionette_L":  (cx - Emu(int(Inches(0.45).emu*scale)), cy + Emu(int(Inches(0.95).emu*scale))),
        "marionette_R":  (cx + Emu(int(Inches(0.45).emu*scale)), cy + Emu(int(Inches(0.95).emu*scale))),
        "chin":          (cx, cy + Emu(int(Inches(1.25).emu*scale))),
        "jaw_L":         (cx - Emu(int(Inches(1.3).emu*scale)), cy + Emu(int(Inches(1.0).emu*scale))),
        "jaw_R":         (cx + Emu(int(Inches(1.3).emu*scale)), cy + Emu(int(Inches(1.0).emu*scale))),
    }


# ─────────────────────────────────────────────────────────────────────────────
# Slide chrome
# ─────────────────────────────────────────────────────────────────────────────

def add_header_band(slide, title_he: str, title_en: str = "", subtitle_he: str = "",
                    accent_color: RGBColor = GOLD):
    # Top gold bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.18), fill=accent_color)
    # Title block on the right
    add_textbox(slide, Inches(6.5), Inches(0.35), Inches(6.6), Inches(0.7),
                title_he, font=FONT_HEB, size=28, bold=True, color=INK,
                align="right", rtl=True)
    if title_en:
        add_textbox(slide, Inches(6.5), Inches(1.05), Inches(6.6), Inches(0.4),
                    title_en, font=FONT_LAT, size=14, bold=False, color=GOLD_DARK,
                    align="right", rtl=False, italic=True)
    if subtitle_he:
        add_textbox(slide, Inches(0.3), Inches(0.55), Inches(6.0), Inches(0.5),
                    subtitle_he, font=FONT_HEB, size=14, bold=False, color=INK_SOFT,
                    align="left", rtl=True)
    # Thin separator under header
    add_line(slide, Inches(0.4), Inches(1.5), Inches(12.93), Inches(1.5),
             color=LINE_SOFT, weight=0.75)


def add_footer(slide, page_num: int, total: int, section: str = ""):
    # Bottom divider
    add_line(slide, Inches(0.4), Inches(7.05), Inches(12.93), Inches(7.05),
             color=LINE_SOFT, weight=0.5)
    # Left: brand
    add_textbox(slide, Inches(0.4), Inches(7.1), Inches(5.0), Inches(0.35),
                "AAA Academy  ·  Restylane Clinical Reference",
                font=FONT_LAT, size=9, color=TAUPE, align="left", rtl=False)
    # Center: section
    if section:
        add_textbox(slide, Inches(5.4), Inches(7.1), Inches(2.5), Inches(0.35),
                    section, font=FONT_HEB, size=9, color=TAUPE,
                    align="center", rtl=True)
    # Right: page
    add_textbox(slide, Inches(11.5), Inches(7.1), Inches(1.43), Inches(0.35),
                f"{page_num} / {total}", font=FONT_LAT, size=9, color=TAUPE,
                align="right", rtl=False)


def add_aaa_badge(slide, left=Inches(0.4), top=Inches(0.35)):
    """Compact AAA Academy badge in the corner."""
    add_rect(slide, left, top, Inches(1.6), Inches(0.7),
             fill=WHITE_WARM, line=GOLD, line_w=0.75,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, left, top + Inches(0.06), Inches(1.6), Inches(0.3),
                "AAA", font=FONT_LAT, size=18, bold=True, color=GOLD_DARK,
                align="center", rtl=False)
    add_textbox(slide, left, top + Inches(0.36), Inches(1.6), Inches(0.28),
                "Academy", font=FONT_LAT, size=9, color=TAUPE,
                align="center", rtl=False, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders (each returns nothing; appends to the presentation)
# ─────────────────────────────────────────────────────────────────────────────

PRODUCTS = {
    # NASHA family
    "classic": {
        "name_en": "Restylane (Classic)",
        "name_he": "רסטילן קלאסי",
        "family": "NASHA",
        "color": NASHA_BLUE,
        "color_lt": NASHA_BLUE_LT,
        "g_prime": "בינוני-גבוה",
        "indication_he": "קמטים שטחיים עד בינוניים, שפתיים עדינות",
        "depth": "דרמה בינונית",
        "longevity_months": "6–9",
        "technique": "מחט / קנולה דקה",
        "tagline_he": "הקלאסיקה - יציבות וקווי מתאר עדינים",
    },
    "lyft": {
        "name_en": "Restylane Lyft",
        "name_he": "רסטילן ליפט",
        "family": "NASHA",
        "color": NASHA_BLUE,
        "color_lt": NASHA_BLUE_LT,
        "g_prime": "גבוה",
        "indication_he": "עצמות לחיים, זווית לסת, חידוד סנטר, גב כפות ידיים",
        "depth": "סופרא-פריאוסטאלי (עמוק על העצם)",
        "longevity_months": "12–18",
        "technique": "Bolus עמוק, מחט/קנולה",
        "tagline_he": "Bone Mimicry – הרמה והגדרה מבנית",
    },
    "eyelight": {
        "name_en": "Restylane Eyelight",
        "name_he": "רסטילן אייליט",
        "family": "NASHA",
        "color": NASHA_BLUE,
        "color_lt": NASHA_BLUE_LT,
        "g_prime": "בינוני-גבוה (נמוך-נפיחות)",
        "indication_he": "Tear-Trough (שקע מתחת לעין) – אזור עדין במיוחד",
        "depth": "פריאוסטאלי / עמוק מעל העצם",
        "longevity_months": "9–12",
        "technique": "קנולה בלבד, Aliquots קטנים",
        "tagline_he": "פתרון ייעודי לאזור התת-עיני",
    },
    "vital": {
        "name_en": "Restylane Skinboosters Vital",
        "name_he": "רסטילן ויטל",
        "family": "NASHA",
        "color": NASHA_BLUE,
        "color_lt": NASHA_BLUE_LT,
        "g_prime": "נמוך (מיועד לסקין-בוסטר)",
        "indication_he": "שיפור טקסטורה, הידרציה, גמישות עור – פנים, צוואר, מחשוף, ידיים",
        "depth": "דרמיס שטחי-בינוני, micro-droplets",
        "longevity_months": "פרוטוקול 3 טיפולים + תחזוקה כל 6 חודשים",
        "technique": "Multi-puncture / מחט מיקרו / NCTF",
        "tagline_he": "Biorevitalization – איכות עור מבפנים",
    },
    # OBT / XpresHAn family
    "refyne": {
        "name_en": "Restylane Refyne",
        "name_he": "רסטילן ריפיין",
        "family": "OBT",
        "color": OBT_PURPLE,
        "color_lt": OBT_PURPLE_LT,
        "g_prime": "נמוך-בינוני",
        "indication_he": "קמטים דינמיים עדינים, NLF שטחי, פרי-אורלי, marionette עדין",
        "depth": "דרמה בינונית",
        "longevity_months": "6–12",
        "technique": "Linear threading / fanning, מחט/קנולה",
        "tagline_he": "Soft & Flexible – טבעיות בתנועה",
    },
    "defyne": {
        "name_en": "Restylane Defyne",
        "name_he": "רסטילן דיפיין",
        "family": "OBT",
        "color": OBT_PURPLE,
        "color_lt": OBT_PURPLE_LT,
        "g_prime": "בינוני-גבוה (גמיש)",
        "indication_he": "קמטי NLF עמוקים יותר, Marionette, עיצוב סנטר, Pre-jowl sulcus",
        "depth": "דרמה עמוקה / תת-עורי שטחי",
        "longevity_months": "12 חודשים בממוצע",
        "technique": "Bolus + linear, קנולה מומלצת בסנטר",
        "tagline_he": "Tissue Mimicry – תמיכה דינמית באזורי תנועה",
    },
    "volyme": {
        "name_en": "Restylane Volyme",
        "name_he": "רסטילן וולים",
        "family": "OBT",
        "color": OBT_PURPLE,
        "color_lt": OBT_PURPLE_LT,
        "g_prime": "בינוני-גבוה, גמיש",
        "indication_he": "החזרת נפח רך באזור הלחי, פיזור הרמוני, מראה ׳בוטן׳ בלחי",
        "depth": "תת-עורי / interfascial (Sub-SMAS)",
        "longevity_months": "12–18",
        "technique": "Fan/Interfascial עם קנולה ארוכה",
        "tagline_he": "Soft Volume – נפח טבעי, תנועה הרמונית",
    },
    "kysse": {
        "name_en": "Restylane Kysse",
        "name_he": "רסטילן קיס",
        "family": "OBT",
        "color": OBT_PURPLE,
        "color_lt": OBT_PURPLE_LT,
        "g_prime": "בינוני, גמיש מאוד",
        "indication_he": "מילוי שפתיים דינמי, הגדרת קווי שפה, lip-flip",
        "depth": "Submucosal / vermillion border",
        "longevity_months": "עד 12 חודשים (צבע עד 8)",
        "technique": "Linear retrograde + tenting, מחט עדינה",
        "tagline_he": "Kissable Lips – נוחות, צבע ותנועה",
    },
    "vital_light": {
        "name_en": "Restylane Vital Light",
        "name_he": "רסטילן ויטל לייט",
        "family": "OBT",
        "color": OBT_PURPLE,
        "color_lt": OBT_PURPLE_LT,
        "g_prime": "נמוך מאוד",
        "indication_he": "עור דק, גוון בהיר, צוואר עדין, אזור פרי-אוקולרי",
        "depth": "דרמיס שטחי",
        "longevity_months": "פרוטוקול 3 טיפולים + תחזוקה כל 6 חודשים",
        "technique": "Micro-droplets, מחט 30G",
        "tagline_he": "Lightweight Hydration – לעור עדין במיוחד",
    },
}


def build_slide_layout(prs):
    """Use blank layout for full design control."""
    return prs.slide_layouts[6]  # blank


# ─── Slide 1: Title ──────────────────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_gradient(slide, CREAM, CREAM_DEEP)

    # Decorative gold bands
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.4), fill=GOLD)
    add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.4), fill=GOLD_DARK)

    # Brand chip
    add_aaa_badge(slide, left=Inches(0.6), top=Inches(0.7))

    # Restylane logotype hint (text-based)
    add_textbox(slide, Inches(10.0), Inches(0.85), Inches(2.9), Inches(0.6),
                "RESTYLANE®", font=FONT_LAT, size=22, bold=True,
                color=GOLD_DARK, align="right", rtl=False)
    add_textbox(slide, Inches(10.0), Inches(1.35), Inches(2.9), Inches(0.4),
                "Hyaluronic Acid Fillers", font=FONT_LAT, size=11,
                color=TAUPE, align="right", rtl=False, italic=True)

    # Hero title
    add_textbox(slide, Inches(1.0), Inches(2.6), Inches(11.3), Inches(1.4),
                "ליין Restylane:", font=FONT_HEB, size=58, bold=True,
                color=INK, align="center", rtl=True)
    add_textbox(slide, Inches(1.0), Inches(3.8), Inches(11.3), Inches(1.1),
                "מדריך קליני מקיף לרופאים", font=FONT_HEB, size=42,
                color=GOLD_DARK, align="center", rtl=True)

    # English subtitle
    add_textbox(slide, Inches(1.0), Inches(5.1), Inches(11.3), Inches(0.5),
                "A Complete Clinical Reference — Material, Indication, Outcome",
                font=FONT_LAT, size=18, color=INK_SOFT,
                align="center", rtl=False, italic=True)

    # Decorative divider line
    add_line(slide, Inches(5.5), Inches(5.85), Inches(7.83), Inches(5.85),
             color=GOLD, weight=2.0)

    # Footer credits
    add_textbox(slide, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.4),
                "מצגת מקצועית · לרופאים בלבד · על בסיס מידע יצרן ומאמרים מאושרים",
                font=FONT_HEB, size=14, color=TAUPE, align="center", rtl=True)


# ─── Slide 2: Agenda ─────────────────────────────────────────────────────────

def slide_agenda(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "תוכן עניינים", "Agenda")

    items = [
        ("01", "מבוא: חומצה היאלורונית באסתטיקה",         "HA in Aesthetic Medicine"),
        ("02", "שתי טכנולוגיות: NASHA ו-OBT/XpresHAn",     "Two Crosslinking Technologies"),
        ("03", "ריאולוגיה: G' Prime, קוהזיביות, נפיחות",   "Rheology Fundamentals"),
        ("04", "משפחת NASHA – מבנה ויציבות",               "The NASHA Family"),
        ("05", "משפחת OBT/XpresHAn – גמישות ותנועה",        "The OBT/XpresHAn Family"),
        ("06", "פרופיל מלא לכל מוצר",                       "Per-Product Deep Dive"),
        ("07", "מפת אזורי הזרקה ובחירת מוצר",              "Face Map & Product Picker"),
        ("08", "טבלאות השוואה: Lyft vs Defyne / Volyme",   "Head-to-Head Comparisons"),
        ("09", "טכניקה: מחט מול קנולה, עומקי הזרקה",       "Needle vs Cannula, Planes"),
        ("10", "Longevity – משך השפעה",                    "Duration Table"),
        ("11", "פרוטוקולים משולבים – full-face plans",     "Combination Protocols"),
        ("12", "Contraindications & Adverse Events",         "Safety"),
        ("13", "Hyaluronidase – פירוק וחילוץ",             "Reversal"),
        ("14", "Case studies",                              "3 Clinical Cases"),
        ("15", "Decision tree & סיכום",                    "Summary"),
    ]
    # Two columns
    col_w = Inches(6.0)
    row_h = Inches(0.42)
    base_top = Inches(1.85)
    for i, (num, he, en) in enumerate(items):
        col = i // 8
        row = i % 8
        left = Inches(0.55) + col * Inches(6.4)
        top = base_top + row * row_h
        # number chip
        add_rect(slide, left, top + Inches(0.05), Inches(0.55), Inches(0.32),
                 fill=GOLD, line=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, left, top + Inches(0.04), Inches(0.55), Inches(0.32),
                    num, font=FONT_LAT, size=12, bold=True, color=WHITE_WARM,
                    align="center", rtl=False)
        # Hebrew title (right-aligned within the column)
        add_textbox(slide, left + Inches(0.7), top, Inches(5.2), Inches(0.32),
                    he, font=FONT_HEB, size=14, bold=True, color=INK,
                    align="right", rtl=True)
        # English subtitle
        add_textbox(slide, left + Inches(0.7), top + Inches(0.22), Inches(5.2),
                    Inches(0.22), en, font=FONT_LAT, size=9, color=TAUPE,
                    align="right", rtl=False, italic=True)

    add_footer(slide, page, total, "תוכן עניינים")


# ─── Slide 3: Intro - HA in aesthetics ───────────────────────────────────────

def slide_intro_ha(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "חומצה היאלורונית באסתטיקה", "Hyaluronic Acid (HA): The Why")

    # Left column: text
    add_textbox(slide, Inches(0.6), Inches(1.85), Inches(7.5), Inches(0.5),
                "מה זה HA?", font=FONT_HEB, size=22, bold=True, color=GOLD_DARK,
                align="right", rtl=True)
    add_textbox(slide, Inches(0.6), Inches(2.35), Inches(7.5), Inches(1.4),
                "פוליסכריד טבעי שנמצא בעור, במפרקים, בנוזל זגוגי - מסוגל לקשור עד פי 1000 ממשקלו במים. עם הגיל הריכוז יורד דרמטית, ואיתו הנפח, הגמישות וההידרציה של העור.",
                font=FONT_HEB, size=15, color=INK_SOFT, align="right", rtl=True,
                line_spacing=1.25)

    add_textbox(slide, Inches(0.6), Inches(3.85), Inches(7.5), Inches(0.5),
                "למה Crosslinking?", font=FONT_HEB, size=22, bold=True, color=GOLD_DARK,
                align="right", rtl=True)
    add_textbox(slide, Inches(0.6), Inches(4.35), Inches(7.5), Inches(1.4),
                "HA חופשי מתפרק תוך 24-48 שעות. קישור צולב (BDDE) הופך אותו ל-Gel יציב שמחזיק חודשים. סוג ה-crosslinking הוא מה שמבדיל בין משפחות מוצרים – וזה שורש כל ההחלטה הקלינית.",
                font=FONT_HEB, size=15, color=INK_SOFT, align="right", rtl=True,
                line_spacing=1.25)

    # Right column: 3 key stats as cards
    stats = [
        ("≈ 50%",    "ירידה בריכוז HA",       "בעור עד גיל 50"),
        ("> 1000×",  "יכולת קשירת מים",        "ביחס למשקל המולקולה"),
        ("BDDE",     "תקן ה-crosslinking",     "Butanediol Diglycidyl Ether"),
    ]
    card_top = Inches(2.0)
    for i, (big, mid, small) in enumerate(stats):
        top = card_top + i * Inches(1.45)
        add_rect(slide, Inches(8.6), top, Inches(4.3), Inches(1.3),
                 fill=WHITE_WARM, line=GOLD, line_w=0.75,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, Inches(8.6), top + Inches(0.1), Inches(4.3), Inches(0.6),
                    big, font=FONT_LAT, size=36, bold=True, color=GOLD_DARK,
                    align="center", rtl=False)
        add_textbox(slide, Inches(8.6), top + Inches(0.7), Inches(4.3), Inches(0.35),
                    mid, font=FONT_HEB, size=14, bold=True, color=INK,
                    align="center", rtl=True)
        add_textbox(slide, Inches(8.6), top + Inches(1.0), Inches(4.3), Inches(0.3),
                    small, font=FONT_LAT, size=10, color=TAUPE,
                    align="center", rtl=False, italic=True)

    # Bottom callout
    add_rect(slide, Inches(0.6), Inches(6.0), Inches(12.13), Inches(0.85),
             fill=SAND, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(0.6), Inches(6.05), Inches(12.13), Inches(0.4),
                "המסר הקליני המרכזי",
                font=FONT_HEB, size=12, bold=True, color=GOLD_DARK,
                align="center", rtl=True)
    add_textbox(slide, Inches(0.6), Inches(6.35), Inches(12.13), Inches(0.4),
                "בחירת המוצר הנכון מתחילה לא מהאזור – אלא מה-Rheology של הג'ל ומהאינדיקציה הקלינית.",
                font=FONT_HEB, size=14, color=INK, align="center", rtl=True)

    add_footer(slide, page, total, "מבוא")


# ─── Slide 4: Two technologies ───────────────────────────────────────────────

def slide_two_technologies(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "שתי טכנולוגיות במשפחה אחת",
                    "NASHA  ·  OBT / XpresHAn", subtitle_he="")

    # Two cards
    col_w = Inches(6.0)
    col_h = Inches(4.8)
    col_top = Inches(1.85)

    # ─ NASHA card
    nasha_l = Inches(0.55)
    add_rect(slide, nasha_l, col_top, col_w, col_h,
             fill=WHITE_WARM, line=NASHA_BLUE, line_w=2.0,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    # blue header strip
    add_rect(slide, nasha_l, col_top, col_w, Inches(0.6),
             fill=NASHA_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, nasha_l, col_top + Inches(0.08), col_w, Inches(0.45),
                "NASHA", font=FONT_LAT, size=24, bold=True, color=WHITE_WARM,
                align="center", rtl=False)
    add_textbox(slide, nasha_l, col_top + Inches(0.7), col_w, Inches(0.4),
                "(Non-Animal Stabilized Hyaluronic Acid)",
                font=FONT_LAT, size=11, color=NASHA_BLUE,
                align="center", rtl=False, italic=True)

    # NASHA particle visual (large grains)
    pic_l = nasha_l + Inches(2.2)
    pic_t = col_top + Inches(1.2)
    for i, (dx, dy, sz) in enumerate([
        (0.0, 0.0, 0.45), (0.55, 0.05, 0.4), (0.25, 0.55, 0.5),
        (0.85, 0.55, 0.35), (0.0, 1.05, 0.42), (0.6, 1.1, 0.45),
        (-0.4, 0.5, 0.38),
    ]):
        add_oval(slide, pic_l + Inches(dx), pic_t + Inches(dy),
                 Inches(sz), Inches(sz),
                 fill=NASHA_BLUE_LT, line=NASHA_BLUE, line_w=0.75)

    # NASHA bullets
    bullets_top = col_top + Inches(3.0)
    nasha_text = (
        "■ Crosslinking נמוך יחסית, חלקיקים גסים יותר\n"
        "■ G' Prime גבוה – יציב, ׳עומד׳ במקום\n"
        "■ דוחף את הרקמה (Bone Mimicry)\n"
        "■ אידיאלי להרמה, יציבות מבנית, חידוד עצם\n"
        "■ פחות נפיחות, פחות תזוזה\n"
        "■ דורש טכניקה ׳כירורגית׳ מדויקת"
    )
    add_textbox(slide, nasha_l + Inches(0.4), bullets_top,
                col_w - Inches(0.8), Inches(1.7),
                nasha_text, font=FONT_HEB, size=13, color=INK,
                align="right", rtl=True, line_spacing=1.3)

    # ─ OBT card
    obt_l = Inches(6.78)
    add_rect(slide, obt_l, col_top, col_w, col_h,
             fill=WHITE_WARM, line=OBT_PURPLE, line_w=2.0,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(slide, obt_l, col_top, col_w, Inches(0.6),
             fill=OBT_PURPLE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, obt_l, col_top + Inches(0.08), col_w, Inches(0.45),
                "OBT / XpresHAn",
                font=FONT_LAT, size=24, bold=True, color=WHITE_WARM,
                align="center", rtl=False)
    add_textbox(slide, obt_l, col_top + Inches(0.7), col_w, Inches(0.4),
                "(Optimal Balance Technology)",
                font=FONT_LAT, size=11, color=OBT_PURPLE,
                align="center", rtl=False, italic=True)

    # OBT mesh visual (smooth grid)
    mesh_l = obt_l + Inches(1.9)
    mesh_t = col_top + Inches(1.2)
    mesh_w = Inches(2.2)
    mesh_h = Inches(1.5)
    rows, cols = 8, 12
    for r in range(rows):
        add_line(slide, mesh_l, mesh_t + Emu(int(mesh_h.emu * r / (rows - 1))),
                 mesh_l + mesh_w, mesh_t + Emu(int(mesh_h.emu * r / (rows - 1))),
                 color=OBT_PURPLE_LT, weight=0.5)
    for c in range(cols):
        add_line(slide, mesh_l + Emu(int(mesh_w.emu * c / (cols - 1))), mesh_t,
                 mesh_l + Emu(int(mesh_w.emu * c / (cols - 1))), mesh_t + mesh_h,
                 color=OBT_PURPLE_LT, weight=0.5)
    # mesh outline
    add_rect(slide, mesh_l, mesh_t, mesh_w, mesh_h, fill=None,
             line=OBT_PURPLE, line_w=1.5)

    bullets_top = col_top + Inches(3.0)
    obt_text = (
        "■ Crosslinking הומוגני, ׳רשת רציפה׳\n"
        "■ G' Prime נמוך-בינוני – Stretch & Recoil\n"
        "■ נע עם הרקמה (Tissue Mimicry)\n"
        "■ אידיאלי לאזורי תנועה ולנפח רך\n"
        "■ נראה טבעי בחיוך ובדיבור\n"
        "■ סלחני יותר טכנית, אך פחות מתאים לעצם"
    )
    add_textbox(slide, obt_l + Inches(0.4), bullets_top,
                col_w - Inches(0.8), Inches(1.7),
                obt_text, font=FONT_HEB, size=13, color=INK,
                align="right", rtl=True, line_spacing=1.3)

    # Bottom rule of thumb
    add_rect(slide, Inches(0.55), Inches(6.75), Inches(12.23), Inches(0.25),
             fill=GOLD, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(0.55), Inches(6.78), Inches(12.23), Inches(0.22),
                "כלל אצבע: NASHA = עצם · OBT = רקמה רכה ותנועה",
                font=FONT_HEB, size=12, bold=True, color=WHITE_WARM,
                align="center", rtl=True)
    add_footer(slide, page, total, "טכנולוגיה")


# ─── Slide 5: Rheology - G' Prime ────────────────────────────────────────────

def slide_rheology(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "ריאולוגיה של ג׳ל HA",
                    "Rheology — G' Prime, Cohesivity, Swelling")

    # Definitions on the right
    defs = [
        ("G' (Elastic Modulus)", "מדד ל׳קשיחות׳ הג'ל – כמה הוא ׳עומד׳ לעומת כמה הוא ׳זורם׳.\nגבוה → תמיכה מבנית, מתאים לעצם.\nנמוך → תנועה טבעית, מתאים לרקמה רכה."),
        ("Cohesivity",          "כמה הג'ל ׳נצמד לעצמו׳ ולא מתפזר. חשוב למניעת migration."),
        ("Swelling",             "כמה הג'ל סופג מים אחרי הזרקה – משפיע על נפיחות מאוחרת."),
        ("Plasticity",          "יכולת לעבור עיצוב אחרי הזרקה (massage / molding)."),
    ]
    top0 = Inches(1.85)
    for i, (term_en, body_he) in enumerate(defs):
        t = top0 + i * Inches(1.15)
        # term chip
        add_rect(slide, Inches(7.5), t, Inches(5.4), Inches(0.4),
                 fill=GOLD_DARK, line=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, Inches(7.5), t + Inches(0.04), Inches(5.4), Inches(0.32),
                    term_en, font=FONT_LAT, size=14, bold=True, color=WHITE_WARM,
                    align="center", rtl=False)
        add_textbox(slide, Inches(7.5), t + Inches(0.45), Inches(5.4), Inches(0.65),
                    body_he, font=FONT_HEB, size=11, color=INK,
                    align="right", rtl=True, line_spacing=1.25)

    # Left: G' Prime spectrum bar (visual)
    bar_l = Inches(0.6)
    bar_t = Inches(2.0)
    bar_w = Inches(6.5)
    bar_h = Inches(0.7)
    # axis label
    add_textbox(slide, bar_l, bar_t - Inches(0.45), bar_w, Inches(0.35),
                "ספקטרום G' Prime  →  קשיחות / יציבות מבנית",
                font=FONT_HEB, size=14, bold=True, color=INK,
                align="center", rtl=True)
    # gradient-look bar built from cells
    cells = 10
    cell_w = Emu(bar_w.emu // cells)
    for i in range(cells):
        # interpolate purple→blue
        t = i / (cells - 1)
        r = int(OBT_PURPLE_LT[0] * (1 - t) + NASHA_BLUE[0] * t)
        g = int(OBT_PURPLE_LT[1] * (1 - t) + NASHA_BLUE[1] * t)
        b = int(OBT_PURPLE_LT[2] * (1 - t) + NASHA_BLUE[2] * t)
        add_rect(slide, bar_l + i * cell_w, bar_t, cell_w, bar_h,
                 fill=RGBColor(r, g, b), line=None)
    # ticks/labels
    products_on_axis = [
        ("Vital",      0.04, OBT_PURPLE),
        ("Refyne",     0.18, OBT_PURPLE),
        ("Kysse",      0.32, OBT_PURPLE),
        ("Classic",    0.46, NASHA_BLUE),
        ("Defyne",     0.55, OBT_PURPLE),
        ("Volyme",     0.65, OBT_PURPLE),
        ("Lyft",       0.92, NASHA_BLUE),
    ]
    for name, pos, color in products_on_axis:
        x = bar_l + Emu(int(bar_w.emu * pos))
        # marker
        add_line(slide, x, bar_t + bar_h, x, bar_t + bar_h + Inches(0.2),
                 color=INK, weight=1.5)
        # label
        add_textbox(slide, x - Inches(0.6), bar_t + bar_h + Inches(0.18),
                    Inches(1.2), Inches(0.3),
                    name, font=FONT_LAT, size=10, bold=True, color=color,
                    align="center", rtl=False)

    # End labels
    add_textbox(slide, bar_l - Inches(0.3), bar_t + bar_h + Inches(0.6),
                Inches(2.5), Inches(0.3),
                "← רך, גמיש (Soft)", font=FONT_HEB, size=11, color=OBT_PURPLE,
                align="left", rtl=True)
    add_textbox(slide, bar_l + bar_w - Inches(2.5), bar_t + bar_h + Inches(0.6),
                Inches(2.8), Inches(0.3),
                "(Firm) קשיח, מבני →", font=FONT_HEB, size=11, color=NASHA_BLUE,
                align="right", rtl=True)

    # Bottom take-away
    add_rect(slide, Inches(0.55), Inches(4.6), Inches(6.6), Inches(2.3),
             fill=WHITE_WARM, line=LINE_SOFT, line_w=0.75,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(0.55), Inches(4.7), Inches(6.6), Inches(0.4),
                "כלל אצבע קליני",
                font=FONT_HEB, size=14, bold=True, color=GOLD_DARK,
                align="center", rtl=True)
    add_textbox(slide, Inches(0.75), Inches(5.05), Inches(6.2), Inches(1.85),
                "G' גבוה  →  עצם, lift, contouring (Lyft)\n"
                "G' בינוני  →  midface volume, סנטר עמוק (Volyme, Defyne)\n"
                "G' נמוך   →  שפתיים, NLF שטחי, אזורי תנועה (Refyne, Kysse)\n"
                "G' מינימלי  →  סקין-בוסטר / hydration (Vital, Vital Light)",
                font=FONT_HEB, size=12, color=INK, align="right", rtl=True,
                line_spacing=1.45)

    add_footer(slide, page, total, "ריאולוגיה")


# ─── Section divider ─────────────────────────────────────────────────────────

def slide_section(prs, page, total, num, title_he, title_en, color, accent_color, body_he=""):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_gradient(slide, CREAM, CREAM_DEEP)

    # Left accent panel
    add_rect(slide, 0, 0, Inches(4.5), SLIDE_H, fill=color)
    add_rect(slide, Inches(4.3), 0, Inches(0.2), SLIDE_H, fill=accent_color)

    # Section number (large)
    add_textbox(slide, Inches(0.4), Inches(1.4), Inches(3.7), Inches(2.0),
                num, font=FONT_LAT, size=180, bold=True, color=WHITE_WARM,
                align="center", rtl=False)
    add_textbox(slide, Inches(0.4), Inches(3.7), Inches(3.7), Inches(0.5),
                "SECTION", font=FONT_LAT, size=14, bold=True,
                color=WHITE_WARM, align="center", rtl=False)

    # Right side: section titles
    add_textbox(slide, Inches(5.0), Inches(2.5), Inches(7.9), Inches(1.2),
                title_he, font=FONT_HEB, size=50, bold=True, color=INK,
                align="right", rtl=True)
    add_textbox(slide, Inches(5.0), Inches(3.7), Inches(7.9), Inches(0.6),
                title_en, font=FONT_LAT, size=22, color=color,
                align="right", rtl=False, italic=True)
    add_line(slide, Inches(11.5), Inches(4.4), Inches(12.9), Inches(4.4),
             color=accent_color, weight=3.0)

    if body_he:
        add_textbox(slide, Inches(5.0), Inches(4.7), Inches(7.9), Inches(1.8),
                    body_he, font=FONT_HEB, size=15, color=INK_SOFT,
                    align="right", rtl=True, line_spacing=1.4)

    add_footer(slide, page, total, title_he)


# ─── Product slide template ──────────────────────────────────────────────────

def slide_product(prs, page, total, key, overlay_kind: str):
    p = PRODUCTS[key]
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, p["name_he"], p["name_en"],
                    accent_color=p["color"])

    # Family chip (top-left)
    fam_color = p["color"]
    add_rect(slide, Inches(0.4), Inches(0.5), Inches(1.7), Inches(0.45),
             fill=fam_color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(0.4), Inches(0.55), Inches(1.7), Inches(0.35),
                p["family"], font=FONT_LAT, size=14, bold=True,
                color=WHITE_WARM, align="center", rtl=False)

    # Tagline under header
    add_textbox(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(0.45),
                p["tagline_he"], font=FONT_HEB, size=18, italic=True,
                color=p["color"], align="center", rtl=True)

    # ─ Left half: clean abstract product visual ──────────────────────────────
    face_panel_l = Inches(0.55)
    face_panel_t = Inches(2.15)
    face_panel_w = Inches(5.8)
    face_panel_h = Inches(4.5)
    add_rect(slide, face_panel_l, face_panel_t, face_panel_w, face_panel_h,
             fill=WHITE_WARM, line=p["color_lt"], line_w=1.0,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

    concept_en = draw_product_visual(slide, face_panel_l, face_panel_t,
                                      face_panel_w, face_panel_h,
                                      overlay_kind, p["color"], p["color_lt"])
    # Caption under the visual
    add_textbox(slide, face_panel_l, face_panel_t + face_panel_h - Inches(0.55),
                face_panel_w, Inches(0.4),
                concept_en, font=FONT_LAT, size=14, italic=True, bold=True,
                color=p["color"], align="center", rtl=False)

    # ─ Right half: data card ────────────────────────────────────────────────
    data_l = Inches(6.7)
    data_t = Inches(2.15)
    data_w = Inches(6.2)
    data_h = Inches(4.5)
    add_rect(slide, data_l, data_t, data_w, data_h,
             fill=WHITE_WARM, line=LINE_SOFT, line_w=0.75,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

    rows = [
        ("מטרה",       p["tagline_he"]),
        ("אינדיקציה",  p["indication_he"]),
        ("עומק הזרקה", p["depth"]),
        ("טכניקה",     p["technique"]),
        ("G' Prime",   p["g_prime"]),
        ("Longevity",  p["longevity_months"]),
    ]
    row_h = Inches(0.7)
    for i, (label, value) in enumerate(rows):
        rt = data_t + Inches(0.2) + i * row_h
        # label chip (right side — RTL: label on right)
        add_rect(slide, data_l + data_w - Inches(1.6), rt,
                 Inches(1.45), Inches(0.5),
                 fill=p["color"], line=None,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, data_l + data_w - Inches(1.6), rt + Inches(0.07),
                    Inches(1.45), Inches(0.36),
                    label, font=FONT_HEB, size=12, bold=True,
                    color=WHITE_WARM, align="center", rtl=True)
        # value
        add_textbox(slide, data_l + Inches(0.2), rt,
                    data_w - Inches(1.9), Inches(0.6),
                    value, font=FONT_HEB, size=12, color=INK,
                    align="right", rtl=True, anchor="middle",
                    line_spacing=1.15)

    add_footer(slide, page, total, p["name_he"])


# ─── Slide: Face Map ─────────────────────────────────────────────────────────

def slide_face_map(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "מפת אזורי הזרקה", "Face Map — Product Picker")

    # Frontal face centered
    cx = Inches(6.66)
    cy = Inches(4.2)
    lm = draw_face_silhouette_front(slide, cx, cy, scale=1.0)

    # First draw markers on both sides for visual symmetry
    sym_markers = [
        ("temple_L", NASHA_BLUE), ("temple_R", NASHA_BLUE),
        ("tear_trough_L", NASHA_BLUE), ("tear_trough_R", NASHA_BLUE),
        ("cheek_L", OBT_PURPLE), ("cheek_R", OBT_PURPLE),
        ("nlf_L", OBT_PURPLE), ("nlf_R", OBT_PURPLE),
        ("marionette_L", OBT_PURPLE), ("marionette_R", OBT_PURPLE),
        ("jaw_L", NASHA_BLUE), ("jaw_R", NASHA_BLUE),
        ("chin", NASHA_BLUE), ("lips", OBT_PURPLE),
    ]
    for lname, color in sym_markers:
        x, y = lm[lname]
        add_oval(slide, x - Inches(0.12), y - Inches(0.12),
                 Inches(0.24), Inches(0.24), fill=color,
                 line=WHITE_PURE, line_w=1.5)

    # Then place labels: 7 on the left (point to L-side landmarks),
    # plus chin/lips at the bottom and right respectively.
    # Each label: (landmark, color, he, en, product, label_x_in, label_y_in,
    #              text_align_in_box)
    left_labels = [
        ("temple_L",     NASHA_BLUE, "רקה",         "Temple",      "Lyft / Volyme",     0.4, 1.0),
        ("tear_trough_L",NASHA_BLUE, "שקע תת-עיני", "Tear-Trough", "Eyelight",          0.4, 1.85),
        ("cheek_L",      OBT_PURPLE, "לחי",         "Midface",     "Volyme / Lyft",     0.4, 2.7),
        ("nlf_L",        OBT_PURPLE, "קפל אף-שפתי","NLF",         "Defyne / Refyne",   0.4, 3.55),
        ("marionette_L", OBT_PURPLE, "Marionette",  "Marionette",  "Defyne",            0.4, 4.4),
        ("jaw_L",        NASHA_BLUE, "זווית לסת",  "Jaw angle",   "Lyft",              0.4, 5.25),
    ]
    right_labels = [
        ("lips",         OBT_PURPLE, "שפתיים",      "Lips",        "Kysse",             10.5, 4.2),
        ("chin",         NASHA_BLUE, "סנטר",        "Chin/Pre-jowl","Lyft + Defyne",    10.5, 5.05),
    ]
    box_w = Inches(2.8)
    box_h = Inches(0.55)
    for entries in (left_labels, right_labels):
        for lname, color, he, en, prod, lx, ly in entries:
            x, y = lm[lname]
            label_x = Inches(lx)
            label_y = Inches(ly)
            # leader line from marker to box's vertical center, then to its near edge
            box_near_x = label_x + (box_w if entries is left_labels else Emu(0))
            box_mid_y = label_y + box_h // 2
            add_line(slide, x, y, box_near_x, box_mid_y,
                     color=color, weight=0.75)
            # box
            add_rect(slide, label_x, label_y, box_w, box_h,
                     fill=WHITE_WARM, line=color, line_w=0.6,
                     shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            add_textbox(slide, label_x, label_y + Inches(0.02),
                        box_w, Inches(0.28),
                        f"{he}  ·  {en}", font=FONT_HEB, size=10, bold=True,
                        color=INK, align="center", rtl=True)
            add_textbox(slide, label_x, label_y + Inches(0.27),
                        box_w, Inches(0.28),
                        prod, font=FONT_LAT, size=10, color=color,
                        align="center", rtl=False, italic=True)

    # Legend (bottom-center, below face)
    leg_l = Inches(5.0)
    leg_t = Inches(6.4)
    add_rect(slide, leg_l, leg_t, Inches(1.5), Inches(0.4),
             fill=NASHA_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, leg_l, leg_t + Inches(0.04), Inches(1.5), Inches(0.32),
                "NASHA",
                font=FONT_HEB, size=11, bold=True, color=WHITE_WARM,
                align="center", rtl=True)
    add_textbox(slide, leg_l + Inches(1.55), leg_t + Inches(0.07), Inches(2.0),
                Inches(0.32), "— מבני / עצם",
                font=FONT_HEB, size=11, color=INK_SOFT,
                align="left", rtl=True)
    leg_l2 = Inches(8.7)
    add_rect(slide, leg_l2, leg_t, Inches(1.5), Inches(0.4),
             fill=OBT_PURPLE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, leg_l2, leg_t + Inches(0.04), Inches(1.5), Inches(0.32),
                "OBT",
                font=FONT_HEB, size=11, bold=True, color=WHITE_WARM,
                align="center", rtl=True)
    add_textbox(slide, leg_l2 + Inches(1.55), leg_t + Inches(0.07),
                Inches(2.5), Inches(0.32), "— רקמה רכה / תנועה",
                font=FONT_HEB, size=11, color=INK_SOFT,
                align="left", rtl=True)

    add_footer(slide, page, total, "מפת אזורים")


# ─── Comparison Slides ───────────────────────────────────────────────────────

def slide_comparison(prs, page, total, title_he, title_en, prod_a, prod_b, rows):
    """rows: list of (attr_he, value_a, value_b)."""
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, title_he, title_en)

    a = PRODUCTS[prod_a]
    b = PRODUCTS[prod_b]

    # Two product header cards
    col_w = Inches(5.5)
    col_top = Inches(1.85)
    # A on right (RTL)
    a_l = Inches(7.4)
    add_rect(slide, a_l, col_top, col_w, Inches(1.2),
             fill=a["color"], shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, a_l, col_top + Inches(0.15), col_w, Inches(0.5),
                a["name_he"], font=FONT_HEB, size=22, bold=True,
                color=WHITE_WARM, align="center", rtl=True)
    add_textbox(slide, a_l, col_top + Inches(0.65), col_w, Inches(0.4),
                a["name_en"], font=FONT_LAT, size=12, color=WHITE_WARM,
                align="center", rtl=False, italic=True)

    # B on left
    b_l = Inches(0.4)
    add_rect(slide, b_l, col_top, col_w, Inches(1.2),
             fill=b["color"], shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, b_l, col_top + Inches(0.15), col_w, Inches(0.5),
                b["name_he"], font=FONT_HEB, size=22, bold=True,
                color=WHITE_WARM, align="center", rtl=True)
    add_textbox(slide, b_l, col_top + Inches(0.65), col_w, Inches(0.4),
                b["name_en"], font=FONT_LAT, size=12, color=WHITE_WARM,
                align="center", rtl=False, italic=True)

    # Comparison rows (3 columns: A | attr | B)
    row_top = Inches(3.3)
    row_h = Inches(0.55)
    attr_col_l = Inches(5.9)
    attr_col_w = Inches(1.5)
    for i, (attr_he, va, vb) in enumerate(rows):
        rt = row_top + i * row_h
        bg = WHITE_WARM if i % 2 == 0 else CREAM_DEEP
        # Whole row background
        add_rect(slide, b_l, rt, Inches(12.5), row_h, fill=bg,
                 line=LINE_SOFT, line_w=0.25)
        # Attr chip (center)
        add_rect(slide, attr_col_l, rt + Inches(0.06), attr_col_w,
                 row_h - Inches(0.12), fill=GOLD_DARK,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, attr_col_l, rt + Inches(0.1), attr_col_w,
                    row_h - Inches(0.2),
                    attr_he, font=FONT_HEB, size=11, bold=True,
                    color=WHITE_WARM, align="center", rtl=True, anchor="middle")
        # Value A (right)
        add_textbox(slide, a_l + Inches(0.2), rt + Inches(0.05),
                    col_w - Inches(0.4), row_h - Inches(0.1),
                    va, font=FONT_HEB, size=12, color=INK,
                    align="right", rtl=True, anchor="middle", line_spacing=1.1)
        # Value B (left)
        add_textbox(slide, b_l + Inches(0.2), rt + Inches(0.05),
                    col_w - Inches(0.4), row_h - Inches(0.1),
                    vb, font=FONT_HEB, size=12, color=INK,
                    align="right", rtl=True, anchor="middle", line_spacing=1.1)

    add_footer(slide, page, total, title_he)


# ─── Slide: Techniques (needle vs cannula) ───────────────────────────────────

def slide_technique(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "טכניקה ועומקי הזרקה",
                    "Needle vs Cannula  ·  Injection Planes")

    # Two cards
    col_top = Inches(1.85)
    col_h = Inches(2.8)

    # Needle card
    add_rect(slide, Inches(0.4), col_top, Inches(6.2), col_h,
             fill=WHITE_WARM, line=NASHA_BLUE, line_w=1.5,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(slide, Inches(0.4), col_top, Inches(6.2), Inches(0.5),
             fill=NASHA_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(0.4), col_top + Inches(0.05), Inches(6.2),
                Inches(0.4),
                "מחט  ·  Needle", font=FONT_HEB, size=18, bold=True,
                color=WHITE_WARM, align="center", rtl=True)
    needle_text = (
        "+ דיוק מקסימלי במיקום הג'ל\n"
        "+ אידיאלי לבולוס עמוק על העצם (Lyft)\n"
        "+ Threading עדין בקמטים שטחיים\n"
        "− סיכון Vascular גבוה יותר\n"
        "− יותר חבלה / hematoma\n"
        "− דורש לימוד אנטומיה מעמיק"
    )
    add_textbox(slide, Inches(0.7), col_top + Inches(0.65), Inches(5.6),
                col_h - Inches(0.75),
                needle_text, font=FONT_HEB, size=12, color=INK,
                align="right", rtl=True, line_spacing=1.4)

    # Cannula card
    add_rect(slide, Inches(6.83), col_top, Inches(6.1), col_h,
             fill=WHITE_WARM, line=OBT_PURPLE, line_w=1.5,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(slide, Inches(6.83), col_top, Inches(6.1), Inches(0.5),
             fill=OBT_PURPLE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(6.83), col_top + Inches(0.05), Inches(6.1),
                Inches(0.4),
                "קנולה  ·  Cannula (Blunt-tip)", font=FONT_HEB, size=18,
                bold=True, color=WHITE_WARM, align="center", rtl=True)
    cann_text = (
        "+ סיכון Vascular נמוך משמעותית\n"
        "+ פחות חבלה, יותר נוחות למטופל\n"
        "+ אידיאלי לאזורים מסוכנים: שקע תת-עיני, NLF, רקות\n"
        "+ פיזור נפח גדול עם פחות חורי כניסה\n"
        "− פחות דיוק (לא מתאים לבולוס עצם)\n"
        "− דורש point of entry מתוכנן"
    )
    add_textbox(slide, Inches(7.13), col_top + Inches(0.65), Inches(5.5),
                col_h - Inches(0.75),
                cann_text, font=FONT_HEB, size=12, color=INK,
                align="right", rtl=True, line_spacing=1.4)

    # Planes diagram (bottom half)
    diag_l = Inches(0.4)
    diag_t = Inches(4.95)
    diag_w = Inches(12.5)
    diag_h = Inches(2.0)
    add_rect(slide, diag_l, diag_t, diag_w, diag_h,
             fill=WHITE_WARM, line=LINE_SOFT, line_w=0.5,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

    # Title
    add_textbox(slide, diag_l, diag_t + Inches(0.05), diag_w, Inches(0.35),
                "עומקי הזרקה (Injection Planes) — מהשטח לעומק",
                font=FONT_HEB, size=14, bold=True, color=GOLD_DARK,
                align="center", rtl=True)

    # Horizontal layers
    layers = [
        ("Epidermis",                RGBColor(0xE8, 0xCF, 0xB5)),
        ("Dermis (superficial)",      RGBColor(0xD3, 0xA9, 0x7F)),
        ("Dermis (deep)",             RGBColor(0xB8, 0x80, 0x55)),
        ("Subcutaneous (Fat)",        RGBColor(0xF8, 0xE5, 0xC2)),
        ("Sub-SMAS / Interfascial",   RGBColor(0xC2, 0xA0, 0x70)),
        ("Supraperiosteal (Bone)",    RGBColor(0xE0, 0xD5, 0xC0)),
    ]
    plane_l = diag_l + Inches(2.8)
    plane_t = diag_t + Inches(0.5)
    plane_w = Inches(3.5)
    layer_h = Inches(0.22)
    for i, (name, color) in enumerate(layers):
        ly = plane_t + i * layer_h
        add_rect(slide, plane_l, ly, plane_w, layer_h, fill=color,
                 line=WHITE_PURE, line_w=0.5)
        add_textbox(slide, plane_l - Inches(2.55), ly - Inches(0.02),
                    Inches(2.45), Inches(0.25),
                    name, font=FONT_LAT, size=10, color=INK,
                    align="right", rtl=False, anchor="middle")

    # Mapping: which products go in which plane (right side)
    mapping = [
        ("דרמה שטחית: Vital Light, Vital",          RGBColor(0x6B, 0x4F, 0x3A)),
        ("דרמה בינונית: Classic, Refyne",            OBT_PURPLE),
        ("דרמה עמוקה / SubQ שטחי: Defyne",           OBT_PURPLE),
        ("Sub-SMAS / Interfascial: Volyme",          OBT_PURPLE),
        ("Supraperiosteal: Lyft, Eyelight",          NASHA_BLUE),
        ("Submucosal (שפתיים): Kysse",                OBT_PURPLE),
    ]
    map_l = plane_l + plane_w + Inches(0.6)
    map_t = plane_t - Inches(0.05)
    for i, (txt, color) in enumerate(mapping):
        add_oval(slide, map_l, map_t + i * Inches(0.22) + Inches(0.05),
                 Inches(0.12), Inches(0.12), fill=color, line=None)
        add_textbox(slide, map_l + Inches(0.2),
                    map_t + i * Inches(0.22) - Inches(0.03),
                    Inches(4.5), Inches(0.28),
                    txt, font=FONT_HEB, size=10, color=INK,
                    align="right", rtl=True, anchor="middle")

    add_footer(slide, page, total, "טכניקה")


# ─── Slide: Longevity table ──────────────────────────────────────────────────

def slide_longevity(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "Longevity — משך השפעה",
                    "Average Duration by Product")

    # Each row: (name, family, label, start_months, end_months, color)
    rows = [
        ("Vital / Vital Light", "OBT/NASHA", "פרוטוקול × 3",  1,  4,  RGBColor(0xA8, 0xBE, 0xD8)),
        ("Restylane Classic",   "NASHA",     "6–9",            6,  9,  NASHA_BLUE_LT),
        ("Restylane Refyne",    "OBT",       "6–12",           6, 12,  OBT_PURPLE_LT),
        ("Restylane Kysse",     "OBT",       "8–12",           8, 12,  OBT_PURPLE_LT),
        ("Restylane Eyelight",  "NASHA",     "9–12",           9, 12,  NASHA_BLUE_LT),
        ("Restylane Defyne",    "OBT",       "10–12",         10, 12,  OBT_PURPLE),
        ("Restylane Volyme",    "OBT",       "12–18",         12, 18,  OBT_PURPLE),
        ("Restylane Lyft",      "NASHA",     "12–18",         12, 18,  NASHA_BLUE),
    ]

    # Axis
    axis_top = Inches(2.0)
    axis_left = Inches(4.5)
    axis_w = Inches(7.8)
    row_h = Inches(0.55)
    max_months = 18

    # Header row
    add_textbox(slide, Inches(0.4), axis_top - Inches(0.4), Inches(4.0),
                Inches(0.3), "מוצר  ·  Product",
                font=FONT_HEB, size=12, bold=True, color=GOLD_DARK,
                align="right", rtl=True)
    add_textbox(slide, axis_left, axis_top - Inches(0.4), axis_w,
                Inches(0.3),
                "0 ────────── 6 ──────────── 12 ─────────── 18 חודשים",
                font=FONT_LAT, size=11, color=GOLD_DARK,
                align="left", rtl=False)

    # 6/12/18 month vertical markers
    for m in (6, 12, 18):
        x = axis_left + Emu(int(axis_w.emu * m / max_months))
        add_line(slide, x, axis_top, x, axis_top + Inches(len(rows) * 0.55) + Inches(0.2),
                 color=LINE_SOFT, weight=0.5)

    for i, (name, fam, label, start, end, color) in enumerate(rows):
        rt = axis_top + i * row_h
        # alternating row bg
        bg = WHITE_WARM if i % 2 == 0 else CREAM_DEEP
        add_rect(slide, Inches(0.4), rt, Inches(12.5), row_h, fill=bg,
                 line=None)
        # product name (right)
        add_textbox(slide, Inches(0.4), rt + Inches(0.07), Inches(3.0),
                    Inches(0.3),
                    name, font=FONT_HEB, size=12, bold=True, color=INK,
                    align="right", rtl=True)
        add_textbox(slide, Inches(0.4), rt + Inches(0.3), Inches(3.0),
                    Inches(0.22),
                    fam, font=FONT_LAT, size=9, color=TAUPE,
                    align="right", rtl=False, italic=True)
        # Duration bar
        bar_x = axis_left + Emu(int(axis_w.emu * start / max_months))
        bar_w = Emu(int(axis_w.emu * (end - start) / max_months))
        add_rect(slide, bar_x, rt + Inches(0.13), bar_w,
                 row_h - Inches(0.26), fill=color, line=None,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        # label inside bar
        add_textbox(slide, bar_x, rt + Inches(0.12), bar_w,
                    row_h - Inches(0.24),
                    label, font=FONT_HEB, size=10, bold=True, color=INK,
                    align="center", rtl=True, anchor="middle")

    # Footnote
    add_textbox(slide, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.4),
                "* טווחים מתבססים על מידע יצרן ומחקרים קליניים. משך השפעה תלוי באזור, בכמות, במטבוליזם של המטופל ובטכניקה.",
                font=FONT_HEB, size=10, color=TAUPE, italic=True,
                align="right", rtl=True)
    add_footer(slide, page, total, "Longevity")


# ─── Slide: Combination protocols ────────────────────────────────────────────

def slide_combinations(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "פרוטוקולים משולבים",
                    "Full-Face Combination Protocols")

    plans = [
        {
            "title": "אישה 40–50, פנים יבשות וירידת נפח קלה",
            "subtitle": "Mid-aging skin quality + early volume loss",
            "color": OBT_PURPLE,
            "steps": [
                "1) Skinboosters Vital × 3 (חודש בין טיפול) — איכות עור כללית",
                "2) Volyme 1ml בכל לחי — החזרת נפח midface",
                "3) Defyne 0.5ml ב-NLF + Marionette — תמיכה דינמית",
                "4) Kysse 0.5–1ml — מילוי שפתיים עדין",
            ],
        },
        {
            "title": "אישה 55+, ירידת נפח בולטת ו-jowls מתפתחים",
            "subtitle": "Mature face — structural lift + soft volume",
            "color": NASHA_BLUE,
            "steps": [
                "1) Lyft 1ml zygoma + temple — הרמה מבנית עמוקה (סופרא-פריאוסטאלי)",
                "2) Volyme 1ml interfascial — נפח לחי הרמוני",
                "3) Lyft 0.5ml zone C (זווית לסת) — חידוד jawline",
                "4) Defyne 0.5–1ml סנטר + pre-jowl — גישור על ה-sulcus",
            ],
        },
        {
            "title": "גבר 35–50, חיזוק קווי מתאר וסנטר חזק",
            "subtitle": "Male contouring — chin & jawline",
            "color": GOLD_DARK,
            "steps": [
                "1) Lyft 1–2ml סנטר — הקרנה קדמית ויציבות",
                "2) Lyft 0.5–1ml × 2 זווית לסת — חידוד קווי מתאר",
                "3) Defyne 0.5ml pre-jowl sulcus — גישור והחלקה",
                "4) Vital בעור — תיקון איכות עור פוסט-גילוח",
            ],
        },
    ]

    col_w = Inches(4.15)
    col_top = Inches(1.85)
    col_h = Inches(4.95)
    gap = Inches(0.1)
    for i, plan in enumerate(plans):
        l = Inches(0.4) + i * (col_w + gap)
        add_rect(slide, l, col_top, col_w, col_h,
                 fill=WHITE_WARM, line=plan["color"], line_w=1.5,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        # header strip
        add_rect(slide, l, col_top, col_w, Inches(1.0),
                 fill=plan["color"], shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, l + Inches(0.2), col_top + Inches(0.1),
                    col_w - Inches(0.4), Inches(0.5),
                    plan["title"], font=FONT_HEB, size=13, bold=True,
                    color=WHITE_WARM, align="right", rtl=True,
                    line_spacing=1.15)
        add_textbox(slide, l + Inches(0.2), col_top + Inches(0.65),
                    col_w - Inches(0.4), Inches(0.35),
                    plan["subtitle"], font=FONT_LAT, size=10,
                    color=WHITE_WARM, align="right", rtl=False, italic=True)
        # Steps
        for j, step in enumerate(plan["steps"]):
            add_textbox(slide, l + Inches(0.2),
                        col_top + Inches(1.15) + j * Inches(0.85),
                        col_w - Inches(0.4), Inches(0.78),
                        step, font=FONT_HEB, size=11, color=INK,
                        align="right", rtl=True, line_spacing=1.3)

    add_footer(slide, page, total, "פרוטוקולים")


# ─── Slide: Pre-treatment consultation ───────────────────────────────────────

def slide_pre_treatment(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "הערכה לפני טיפול",
                    "Pre-Treatment Consultation Checklist")

    sections = [
        ("היסטוריה רפואית", [
            "אלרגיות (כולל ל-Lidocaine, ל-HA, לחלבון תרנגולת לסטרפטוקוקוס)",
            "מחלות אוטואימוניות פעילות / טיפול ביולוגי",
            "הריון או הנקה (אבסולוטי)",
            "טיפולי נוגדי קרישה — מינון ועיתוי",
            "טיפולים אסתטיים קודמים (מוצר, אזור, תאריך, תוצאה)",
            "Cold sores פעיל / היסטוריה (Herpes labialis)",
        ], NASHA_BLUE),
        ("הערכה אסתטית", [
            "ניתוח MD Codes / חלוקה לאזורים",
            "סימטריה ופרופורציות (חוקי שליש, פי הזהב)",
            "איכות עור (Glogau, Fitzpatrick)",
            "הערכה דינמית — חיוך, דיבור, ביטוי",
            "צילום סטנדרטי (קדמי, פרופיל 45°, פרופיל 90°)",
            "ציפיות מטופל מול אפשרי קלינית",
        ], OBT_PURPLE),
        ("Informed Consent", [
            "הסבר על המוצר, מקור החומר, תוקף",
            "סיכונים: hematoma, nodules, infection, vascular event",
            "תוצאות צפויות וזמן ל-onset (24–72 שעות עד 2 שבועות)",
            "Plan ל-Touch up וזמן (כ-2 שבועות)",
            "עלות, פולואפ, אחריות",
            "חתימה דיגיטלית / פיזית בתיק",
        ], GOLD_DARK),
    ]

    col_w = Inches(4.15)
    col_top = Inches(1.85)
    col_h = Inches(4.95)
    gap = Inches(0.1)
    for i, (title, items, color) in enumerate(sections):
        l = Inches(0.4) + i * (col_w + gap)
        add_rect(slide, l, col_top, col_w, col_h,
                 fill=WHITE_WARM, line=color, line_w=1.5,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_rect(slide, l, col_top, col_w, Inches(0.55),
                 fill=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, l, col_top + Inches(0.08), col_w, Inches(0.4),
                    title, font=FONT_HEB, size=16, bold=True,
                    color=WHITE_WARM, align="center", rtl=True)
        for j, item in enumerate(items):
            it_top = col_top + Inches(0.7) + j * Inches(0.7)
            # checkbox
            add_rect(slide, l + Inches(0.2), it_top + Inches(0.05),
                     Inches(0.22), Inches(0.22), fill=WHITE_PURE,
                     line=color, line_w=1.0,
                     shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            add_textbox(slide, l + Inches(0.5), it_top,
                        col_w - Inches(0.7), Inches(0.65),
                        item, font=FONT_HEB, size=11, color=INK,
                        align="right", rtl=True, line_spacing=1.3)

    add_footer(slide, page, total, "Consultation")


# ─── Slide: Contraindications & Safety ───────────────────────────────────────

def slide_contraindications(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "התוויות נגד ובטיחות",
                    "Contraindications & Safety")

    # Two columns: absolute (red) | relative (yellow)
    col_top = Inches(1.85)
    col_h = Inches(3.5)
    col_w = Inches(6.1)

    # Absolute
    add_rect(slide, Inches(0.4), col_top, col_w, col_h,
             fill=WHITE_WARM, line=DANGER, line_w=2.0,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(slide, Inches(0.4), col_top, col_w, Inches(0.55),
             fill=DANGER, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(0.4), col_top + Inches(0.08), col_w, Inches(0.4),
                "אבסולוטי  ·  Absolute",
                font=FONT_HEB, size=18, bold=True, color=WHITE_WARM,
                align="center", rtl=True)
    abs_items = [
        "אלרגיה ידועה ל-HA או ל-BDDE",
        "הריון, הנקה (חוסר מידע)",
        "זיהום עורי פעיל באזור ההזרקה",
        "Cellulitis / Abscess פעיל",
        "Anaphylaxis היסטורי לטיפול דומה",
        "מחלות אוטואימוניות לא מאוזנות",
    ]
    for j, it in enumerate(abs_items):
        it_t = col_top + Inches(0.75) + j * Inches(0.42)
        add_oval(slide, Inches(0.6), it_t + Inches(0.1),
                 Inches(0.12), Inches(0.12), fill=DANGER, line=None)
        add_textbox(slide, Inches(0.85), it_t, col_w - Inches(0.6),
                    Inches(0.35),
                    it, font=FONT_HEB, size=12, color=INK,
                    align="right", rtl=True)

    # Relative
    rl = Inches(6.7)
    add_rect(slide, rl, col_top, col_w, col_h,
             fill=WHITE_WARM, line=WARN, line_w=2.0,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_rect(slide, rl, col_top, col_w, Inches(0.55),
             fill=WARN, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, rl, col_top + Inches(0.08), col_w, Inches(0.4),
                "יחסי  ·  Relative",
                font=FONT_HEB, size=18, bold=True, color=WHITE_WARM,
                align="center", rtl=True)
    rel_items = [
        "טיפול נוגדי קרישה / NSAIDs בשבוע האחרון",
        "Herpes labialis פעיל / היסטוריה (פרופ׳ Acyclovir)",
        "פרוצדורות לייזר / Peelings טריים (פחות מ-2 שבועות)",
        "Fillers שאינם HA באותו אזור (Sculptra, PMMA)",
        "מטופל לא משתף פעולה / ציפיות בלתי-ריאליות",
        "מצב חיסוני מוחלש זמני",
    ]
    for j, it in enumerate(rel_items):
        it_t = col_top + Inches(0.75) + j * Inches(0.42)
        add_oval(slide, rl + Inches(0.2), it_t + Inches(0.1),
                 Inches(0.12), Inches(0.12), fill=WARN, line=None)
        add_textbox(slide, rl + Inches(0.45), it_t,
                    col_w - Inches(0.6), Inches(0.35),
                    it, font=FONT_HEB, size=12, color=INK,
                    align="right", rtl=True)

    # Adverse Events strip
    ae_top = Inches(5.55)
    add_rect(slide, Inches(0.4), ae_top, Inches(12.5), Inches(1.4),
             fill=WHITE_WARM, line=GOLD, line_w=1.0,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(0.4), ae_top + Inches(0.07), Inches(12.5),
                Inches(0.35),
                "Adverse Events — מוקדמים ומאוחרים",
                font=FONT_HEB, size=14, bold=True, color=GOLD_DARK,
                align="center", rtl=True)
    add_textbox(slide, Inches(0.6), ae_top + Inches(0.45), Inches(5.8),
                Inches(0.9),
                "מוקדמים: אריתמה, בצקת, hematoma, רגישות, כאב.\n"
                "ניהול: קרח, ראש מורם, Arnica, NSAID לפי הצורך.",
                font=FONT_HEB, size=11, color=INK, align="right", rtl=True,
                line_spacing=1.4)
    add_textbox(slide, Inches(6.8), ae_top + Inches(0.45), Inches(6.0),
                Inches(0.9),
                "מאוחרים: nodules, granuloma, biofilm, Tyndall effect, occlusion vascular (חירום!).\n"
                "ניהול: Hyaluronidase, אנטיביוטיקה, פניה דחופה.",
                font=FONT_HEB, size=11, color=INK, align="right", rtl=True,
                line_spacing=1.4)

    add_footer(slide, page, total, "Safety")


# ─── Slide: Hyaluronidase ────────────────────────────────────────────────────

def slide_hyaluronidase(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "Hyaluronidase — חילוץ ופירוק",
                    "Reversal Protocol")

    # Hero box
    add_rect(slide, Inches(0.4), Inches(1.85), Inches(12.5), Inches(0.7),
             fill=DANGER, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(0.4), Inches(1.9), Inches(12.5), Inches(0.6),
                "Vascular Occlusion = חירום! התחלת Hyaluronidase מיידית, ללא המתנה לאישור.",
                font=FONT_HEB, size=16, bold=True, color=WHITE_WARM,
                align="center", rtl=True)

    # Dosing table
    table_t = Inches(2.85)
    add_rect(slide, Inches(0.4), table_t, Inches(6.2), Inches(3.7),
             fill=WHITE_WARM, line=LINE_SOFT, line_w=0.5,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(0.4), table_t + Inches(0.1), Inches(6.2),
                Inches(0.4),
                "מינון מומלץ (Hylenex, Hyalase)",
                font=FONT_HEB, size=14, bold=True, color=GOLD_DARK,
                align="center", rtl=True)

    dose_rows = [
        ("Nodule קטן",            "50–150 IU מקומית"),
        ("Asymmetry / overcorrection","100–300 IU"),
        ("Vascular occlusion (DANGER!)", "500–1500 IU + flooding כל שעה"),
        ("בליטה לימפטית",          "75–150 IU מקומית, חזרה לפי הצורך"),
        ("Tear-trough Tyndall",   "5–10 IU בלבד, מינון נמוך מאוד"),
    ]
    for i, (case, dose) in enumerate(dose_rows):
        rt = table_t + Inches(0.6) + i * Inches(0.55)
        bg = CREAM_DEEP if i % 2 == 0 else WHITE_WARM
        add_rect(slide, Inches(0.4), rt, Inches(6.2), Inches(0.5), fill=bg)
        add_textbox(slide, Inches(0.6), rt + Inches(0.05), Inches(3.0),
                    Inches(0.4),
                    case, font=FONT_HEB, size=12, color=INK,
                    align="right", rtl=True, anchor="middle")
        add_textbox(slide, Inches(3.6), rt + Inches(0.05), Inches(2.9),
                    Inches(0.4),
                    dose, font=FONT_HEB, size=12, bold=True, color=DANGER,
                    align="right", rtl=True, anchor="middle")

    # Decision protocol
    add_rect(slide, Inches(6.8), table_t, Inches(6.1), Inches(3.7),
             fill=WHITE_WARM, line=LINE_SOFT, line_w=0.5,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(6.8), table_t + Inches(0.1), Inches(6.1),
                Inches(0.4),
                "פרוטוקול חירום — Vascular Event",
                font=FONT_HEB, size=14, bold=True, color=DANGER,
                align="center", rtl=True)
    steps = [
        "1) זיהוי: blanching, כאב חד, livedo, אובדן ראייה",
        "2) הפסקת הזרקה מיידית",
        "3) Hyaluronidase 500–1500 IU מיידית באזור החסום",
        "4) Massage חזק, חום מקומי",
        "5) Aspirin 325 mg, Nitroglycerin paste (לפי הפרוטוקול)",
        "6) חזרה על Hyaluronidase כל שעה עד שיפור perfusion",
        "7) במעורבות עינית — דחוף לבית חולים תוך 90 דק׳",
    ]
    for i, s in enumerate(steps):
        add_textbox(slide, Inches(7.0), table_t + Inches(0.5) + i * Inches(0.45),
                    Inches(5.8), Inches(0.4),
                    s, font=FONT_HEB, size=11, color=INK,
                    align="right", rtl=True, line_spacing=1.25)

    # Footnote
    add_textbox(slide, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.4),
                "* רכוש Hyaluronidase ושמור בקליניקה תמיד. ערוך skin test ב-15IU תת-עורי לפני שימוש ראשון במטופל אלרגן.",
                font=FONT_HEB, size=10, color=TAUPE, italic=True,
                align="right", rtl=True)
    add_footer(slide, page, total, "Reversal")


# ─── Case study slide ────────────────────────────────────────────────────────

def slide_case(prs, page, total, num, title_he, demographics, findings,
               plan, outcome, color):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, f"Case {num}: {title_he}", f"Case Study {num}")

    # Top demographics strip
    add_rect(slide, Inches(0.4), Inches(1.85), Inches(12.5), Inches(0.6),
             fill=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(0.4), Inches(1.9), Inches(12.5), Inches(0.5),
                demographics, font=FONT_HEB, size=14, bold=True,
                color=WHITE_WARM, align="center", rtl=True)

    # Three columns: findings / plan / outcome
    col_top = Inches(2.65)
    col_h = Inches(4.0)
    col_w = Inches(4.1)
    gap = Inches(0.13)

    cards = [
        ("Findings  ·  ממצאים",   findings, NASHA_BLUE),
        ("Treatment Plan  ·  תוכנית", plan, GOLD_DARK),
        ("Outcome  ·  תוצאה",     outcome, SUCCESS),
    ]
    for i, (title, items, cc) in enumerate(cards):
        l = Inches(0.4) + i * (col_w + gap)
        add_rect(slide, l, col_top, col_w, col_h, fill=WHITE_WARM,
                 line=cc, line_w=1.5,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_rect(slide, l, col_top, col_w, Inches(0.55), fill=cc,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, l, col_top + Inches(0.08), col_w, Inches(0.4),
                    title, font=FONT_HEB, size=14, bold=True,
                    color=WHITE_WARM, align="center", rtl=True)
        for j, item in enumerate(items):
            add_textbox(slide, l + Inches(0.25),
                        col_top + Inches(0.7) + j * Inches(0.55),
                        col_w - Inches(0.5), Inches(0.55),
                        f"• {item}", font=FONT_HEB, size=11, color=INK,
                        align="right", rtl=True, line_spacing=1.3)

    add_footer(slide, page, total, f"Case {num}")


# ─── Decision tree slide ─────────────────────────────────────────────────────

def slide_decision_tree(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "Decision Tree — איזה מוצר מתי",
                    "Quick-Reference Decision Aid")

    # Root question
    add_rect(slide, Inches(5.0), Inches(1.85), Inches(3.33), Inches(0.7),
             fill=GOLD_DARK, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(5.0), Inches(1.95), Inches(3.33), Inches(0.5),
                "מה האינדיקציה הראשית?",
                font=FONT_HEB, size=14, bold=True, color=WHITE_WARM,
                align="center", rtl=True)

    # 4 branches
    branches = [
        ("איכות עור / hydration",      "Vital · Vital Light",     OBT_PURPLE,  0),
        ("שפתיים",                     "Kysse",                    OBT_PURPLE,  1),
        ("נפח / contouring",            "→ ראה ענף נפח",           NASHA_BLUE,  2),
        ("קמטים דינמיים",              "→ ראה ענף קמטים",          OBT_PURPLE,  3),
    ]
    base_t = Inches(3.1)
    branch_w = Inches(3.0)
    for i, (q, ans, color, col) in enumerate(branches):
        l = Inches(0.4) + col * Inches(3.25)
        # connector
        add_line(slide, Inches(6.66), Inches(2.55),
                 l + branch_w / 2, base_t, color=color, weight=1.5)
        # question chip
        add_rect(slide, l, base_t, branch_w, Inches(0.55), fill=color,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, l, base_t + Inches(0.08), branch_w, Inches(0.4),
                    q, font=FONT_HEB, size=12, bold=True, color=WHITE_WARM,
                    align="center", rtl=True)
        # answer
        add_rect(slide, l, base_t + Inches(0.7), branch_w, Inches(0.55),
                 fill=WHITE_WARM, line=color, line_w=1.0,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_textbox(slide, l, base_t + Inches(0.78), branch_w, Inches(0.4),
                    ans, font=FONT_LAT, size=12, bold=True, color=color,
                    align="center", rtl=False)

    # Sub-branch: Volume / Contouring
    sb_t = Inches(4.65)
    sb_l = Inches(0.4)
    add_rect(slide, sb_l, sb_t, Inches(6.1), Inches(2.15),
             fill=WHITE_WARM, line=NASHA_BLUE, line_w=1.0,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, sb_l, sb_t + Inches(0.05), Inches(6.1), Inches(0.4),
                "ענף נפח: שיקול עומק ומטרה",
                font=FONT_HEB, size=13, bold=True, color=NASHA_BLUE,
                align="center", rtl=True)
    vol_rows = [
        ("עצם / contouring (לסת, סנטר, רקה)",  "Lyft",      NASHA_BLUE),
        ("Midface volume רך, אינטרפסיאלי",     "Volyme",    OBT_PURPLE),
        ("חידוד עצם + תמיכת רקמה (מתקדם)",     "Lyft + Defyne", GOLD_DARK),
        ("Tear-trough תת-עיני",                "Eyelight",  NASHA_BLUE),
    ]
    for i, (q, ans, color) in enumerate(vol_rows):
        rt = sb_t + Inches(0.5) + i * Inches(0.38)
        add_textbox(slide, sb_l + Inches(0.2), rt, Inches(4.3), Inches(0.3),
                    "•  " + q, font=FONT_HEB, size=11, color=INK,
                    align="right", rtl=True, anchor="middle")
        add_textbox(slide, sb_l + Inches(4.5), rt, Inches(1.5), Inches(0.3),
                    ans, font=FONT_LAT, size=11, bold=True, color=color,
                    align="center", rtl=False, anchor="middle")

    # Sub-branch: Wrinkles
    sb2_l = Inches(6.83)
    add_rect(slide, sb2_l, sb_t, Inches(6.1), Inches(2.15),
             fill=WHITE_WARM, line=OBT_PURPLE, line_w=1.0,
             shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, sb2_l, sb_t + Inches(0.05), Inches(6.1), Inches(0.4),
                "ענף קמטים: עומק ודינמיות",
                font=FONT_HEB, size=13, bold=True, color=OBT_PURPLE,
                align="center", rtl=True)
    wrk_rows = [
        ("קמטים שטחיים, סטטיים",               "Classic",       NASHA_BLUE),
        ("קמטים שטחיים בתנועה",                "Refyne",        OBT_PURPLE),
        ("NLF עמוק + Marionette",              "Defyne",        OBT_PURPLE),
        ("עור דק / פרי-אוקולרי",              "Vital Light",   OBT_PURPLE),
    ]
    for i, (q, ans, color) in enumerate(wrk_rows):
        rt = sb_t + Inches(0.5) + i * Inches(0.38)
        add_textbox(slide, sb2_l + Inches(0.2), rt, Inches(4.3),
                    Inches(0.3),
                    "•  " + q, font=FONT_HEB, size=11, color=INK,
                    align="right", rtl=True, anchor="middle")
        add_textbox(slide, sb2_l + Inches(4.5), rt, Inches(1.5),
                    Inches(0.3),
                    ans, font=FONT_LAT, size=11, bold=True, color=color,
                    align="center", rtl=False, anchor="middle")

    add_footer(slide, page, total, "Decision Tree")


# ─── Summary slide ───────────────────────────────────────────────────────────

def slide_summary(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_solid(slide, CREAM)
    add_header_band(slide, "סיכום — ארבעת המסרים", "Key Take-Aways")

    keys = [
        ("01", "תחילת ההחלטה היא ב-Rheology",
               "בחר משפחה (NASHA/OBT) לפי איכות הג'ל הנדרשת באזור — לא לפי שם המוצר.",
               GOLD_DARK),
        ("02", "NASHA = עצם · OBT = רקמה רכה ותנועה",
               "Bone Mimicry להרמה ולחידוד מבני; Tissue Mimicry לאזורי תנועה ולנפח רך.",
               NASHA_BLUE),
        ("03", "שלב טכניקה עם אנטומיה",
               "קנולה באזורים סכנתיים (NLF, tear-trough, רקות). מחט לבולוס עצם מדויק.",
               OBT_PURPLE),
        ("04", "בטיחות מעל הכל",
               "Hyaluronidase תמיד במלאי. הכר vascular events. תעד ותתעד את הציפיות.",
               DANGER),
    ]

    col_w = Inches(6.1)
    row_h = Inches(2.45)
    for i, (num, title, body, color) in enumerate(keys):
        col = i % 2
        row = i // 2
        l = Inches(0.4) + col * Inches(6.35)
        t = Inches(1.85) + row * row_h
        add_rect(slide, l, t, col_w, row_h - Inches(0.15),
                 fill=WHITE_WARM, line=color, line_w=1.5,
                 shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        # Big number
        add_textbox(slide, l + Inches(0.15), t + Inches(0.15),
                    Inches(1.4), Inches(1.4),
                    num, font=FONT_LAT, size=64, bold=True, color=color,
                    align="center", rtl=False, anchor="middle")
        # Title
        add_textbox(slide, l + Inches(1.55), t + Inches(0.25),
                    col_w - Inches(1.7), Inches(0.6),
                    title, font=FONT_HEB, size=18, bold=True, color=INK,
                    align="right", rtl=True, line_spacing=1.15)
        # Body
        add_textbox(slide, l + Inches(1.55), t + Inches(0.95),
                    col_w - Inches(1.7), Inches(1.2),
                    body, font=FONT_HEB, size=12, color=INK_SOFT,
                    align="right", rtl=True, line_spacing=1.35)

    add_footer(slide, page, total, "סיכום")


# ─── Thank-you / Q&A slide ───────────────────────────────────────────────────

def slide_thankyou(prs, page, total):
    slide = prs.slides.add_slide(build_slide_layout(prs))
    set_slide_bg_gradient(slide, CREAM_DEEP, CREAM)

    # Big thank you
    add_textbox(slide, Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.4),
                "תודה",
                font=FONT_HEB, size=90, bold=True, color=GOLD_DARK,
                align="center", rtl=True)
    add_textbox(slide, Inches(1.0), Inches(3.8), Inches(11.3), Inches(0.6),
                "Q&A  ·  Discussion",
                font=FONT_LAT, size=28, color=INK, align="center",
                rtl=False, italic=True)

    add_line(slide, Inches(5.5), Inches(4.7), Inches(7.83), Inches(4.7),
             color=GOLD, weight=2.5)

    add_textbox(slide, Inches(1.0), Inches(5.0), Inches(11.3), Inches(0.5),
                "תזכורת קלינית: לכל החלטה — Rheology, Indication, Outcome.",
                font=FONT_HEB, size=18, color=INK_SOFT,
                align="center", rtl=True, italic=True)

    # Bottom brand
    add_aaa_badge(slide, left=Inches(5.86), top=Inches(6.1))

    add_textbox(slide, Inches(1.0), Inches(6.95), Inches(11.3), Inches(0.4),
                "AAA Academy · Advanced Aesthetic Anatomy",
                font=FONT_LAT, size=12, color=TAUPE,
                align="center", rtl=False, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Build orchestrator
# ─────────────────────────────────────────────────────────────────────────────

def build(output_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Pre-compute total slide count for footers
    # Title, agenda, intro, two-tech, rheology
    # Sect-A divider + 4 NASHA product slides
    # Sect-B divider + 5 OBT product slides
    # Face map
    # 4 comparison slides
    # Technique, longevity, combinations, pre-treatment, contraindications, hyaluronidase
    # 3 case studies
    # Decision tree, summary, thank-you
    total = (
        5                  # title, agenda, intro, two-tech, rheology
        + 1 + 4            # NASHA divider + 4 products
        + 1 + 5            # OBT divider + 5 products
        + 1                # face map
        + 4                # comparisons
        + 1 + 1 + 1 + 1 + 1 + 1   # tech, long, combo, pretx, contra, hyase
        + 3                # cases
        + 1 + 1 + 1        # tree, summary, thankyou
    )

    page = 1
    slide_title(prs); page += 1
    slide_agenda(prs, page, total); page += 1
    slide_intro_ha(prs, page, total); page += 1
    slide_two_technologies(prs, page, total); page += 1
    slide_rheology(prs, page, total); page += 1

    # NASHA family
    slide_section(prs, page, total, "A",
                  "משפחת NASHA",
                  "Firm · Structural · Bone Mimicry",
                  NASHA_BLUE, GOLD,
                  body_he=(
                      "ג'ל יציב עם G' Prime גבוה, פעולה מבנית.\n"
                      "אידיאלי להרמה, חידוד עצם, ולמיקום עמוק.\n"
                      "ארבעה מוצרים בליין: Classic, Lyft, Eyelight, Vital."
                  )); page += 1
    slide_product(prs, page, total, "classic", "wrinkle_line"); page += 1
    slide_product(prs, page, total, "lyft", "triangle_jaw"); page += 1
    slide_product(prs, page, total, "eyelight", "undereye"); page += 1
    slide_product(prs, page, total, "vital", "dots"); page += 1

    # OBT family
    slide_section(prs, page, total, "B",
                  "משפחת OBT / XpresHAn",
                  "Flexible · Dynamic · Tissue Mimicry",
                  OBT_PURPLE, GOLD,
                  body_he=(
                      "ג'ל הומוגני עם Stretch & Recoil.\n"
                      "אידיאלי לאזורי תנועה ולנפח רך.\n"
                      "חמישה מוצרים בליין: Refyne, Defyne, Volyme, Kysse, Vital Light."
                  )); page += 1
    slide_product(prs, page, total, "refyne", "wrinkle_line"); page += 1
    slide_product(prs, page, total, "defyne", "mesh_jaw"); page += 1
    slide_product(prs, page, total, "volyme", "soft_cheek"); page += 1
    slide_product(prs, page, total, "kysse", "lips"); page += 1
    slide_product(prs, page, total, "vital_light", "dots"); page += 1

    # Face map
    slide_face_map(prs, page, total); page += 1

    # Comparisons
    slide_comparison(prs, page, total,
                     "חידוד ופיסול: Lyft vs Defyne",
                     "Jawline & Chin — Lyft vs Defyne",
                     "lyft", "defyne",
                     rows=[
                         ("משפחה",                "NASHA",                     "OBT/XpresHAn"),
                         ("G' Prime",             "גבוה",                       "בינוני-גבוה (גמיש)"),
                         ("מטרה",                 "יצירת קו עצם וחידוד מבני",  "תמיכה דינמית באזורי תנועה"),
                         ("עומק",                 "סופרא-פריאוסטאלי",          "דרמה עמוקה / SubQ"),
                         ("אינדיקציה",            "זווית לסת, חידוד סנטר",     "סנטר, Pre-jowl, Marionette"),
                         ("טכניקה",               "Bolus עמוק מחט/קנולה",      "Linear + bolus קנולה"),
                         ("Longevity (חודשים)",   "12–18",                      "≈ 12"),
                     ]); page += 1

    slide_comparison(prs, page, total,
                     "עיצוב הלחי: Lyft vs Volyme",
                     "Midface — Lyft vs Volyme",
                     "lyft", "volyme",
                     rows=[
                         ("משפחה",                "NASHA",                     "OBT/XpresHAn"),
                         ("G' Prime",             "גבוה",                       "בינוני-גבוה, גמיש"),
                         ("מטרה",                 "הרמה מבנית מעצם",            "נפח רך, פיזור הרמוני"),
                         ("עומק",                 "סופרא-פריאוסטאלי",          "Interfascial / Sub-SMAS"),
                         ("טכניקה",               "Bolus 'Gunshot' עמוק",      "Fan עם קנולה ארוכה"),
                         ("התווייה",              "ירידה עמוקה, צורך בהרמה",   "ירידת נפח רכה, מראה 'בוטן'"),
                         ("Longevity (חודשים)",   "12–18",                      "12–18"),
                     ]); page += 1

    slide_comparison(prs, page, total,
                     "קמטים דינמיים: Refyne vs Defyne",
                     "Dynamic Wrinkles — Refyne vs Defyne",
                     "refyne", "defyne",
                     rows=[
                         ("משפחה",                "OBT/XpresHAn",              "OBT/XpresHAn"),
                         ("G' Prime",             "נמוך-בינוני",                 "בינוני-גבוה"),
                         ("עומק קמט",             "קמטים שטחיים",                "קמטים בינוניים-עמוקים"),
                         ("אינדיקציה",            "NLF שטחי, פרי-אורלי",       "NLF עמוק, Marionette"),
                         ("עומק הזרקה",           "דרמה בינונית",                "דרמה עמוקה / SubQ"),
                         ("מטרה",                 "טבעיות בתנועה",               "תמיכה דינמית מתקדמת"),
                         ("Longevity (חודשים)",   "6–12",                       "≈ 12"),
                     ]); page += 1

    slide_comparison(prs, page, total,
                     "שפתיים: Kysse vs Classic",
                     "Lips — Kysse vs Classic",
                     "kysse", "classic",
                     rows=[
                         ("משפחה",                "OBT/XpresHAn",              "NASHA"),
                         ("G' Prime",             "בינוני, גמיש מאוד",          "בינוני-גבוה"),
                         ("תנועה / חיוך",         "טבעי, רך בחיוך",             "יותר מורגש בחיוך"),
                         ("צבע ולחות",            "Color retention עד 8 חודשים", "מילוי קלאסי"),
                         ("עומק",                 "Submucosal / vermillion",   "דרמה / submucosal שטחי"),
                         ("טכניקה",               "Linear retrograde עם מחט",  "Linear עם מחט"),
                         ("Longevity (חודשים)",   "עד 12",                      "6–9"),
                     ]); page += 1

    # Practical / safety section
    slide_technique(prs, page, total); page += 1
    slide_longevity(prs, page, total); page += 1
    slide_combinations(prs, page, total); page += 1
    slide_pre_treatment(prs, page, total); page += 1
    slide_contraindications(prs, page, total); page += 1
    slide_hyaluronidase(prs, page, total); page += 1

    # Case studies
    slide_case(prs, page, total, 1,
               "ירידת נפח midface",
               "אישה, גיל 45, Fitzpatrick III, ללא טיפולים קודמים. תלונה: 'עייפות' במרכז הפנים.",
               findings=[
                   "ירידת נפח לחי דו-צדדית",
                   "NLF בולט בעיקר בשמאל",
                   "tear-trough עדין מתחיל",
                   "איכות עור טובה, ללא צלקות",
               ],
               plan=[
                   "Volyme 1ml בכל לחי (interfascial, קנולה)",
                   "Defyne 0.3ml NLF עמוק",
                   "Eyelight 0.2ml tear-trough בקנולה",
                   "המלצה: Vital × 3 לאחר 4 שבועות",
               ],
               outcome=[
                   "החזרת midface volume טבעית",
                   "מילוי NLF עם תנועה טבעית בחיוך",
                   "ללא 'בוטן' או חוסר סימטריה",
                   "מטופלת מרוצה, חוזרת לתחזוקה בעוד 12 חודשים",
               ],
               color=OBT_PURPLE); page += 1

    slide_case(prs, page, total, 2,
               "חידוד לסת וסנטר",
               "גבר, גיל 52, Fitzpatrick II, מבקש 'קו לסת חד' לאחר ירידה במשקל.",
               findings=[
                   "Pre-jowl sulcus דו-צדדי",
                   "זווית לסת קהה",
                   "סנטר ללא הקרנה קדמית",
                   "ללא ירידת נפח שמשמעותית בלחי",
               ],
               plan=[
                   "Lyft 1ml בסנטר (סופרא-פריאוסטאלי, בולוס)",
                   "Lyft 0.7ml × 2 בזווית לסת (Zone C)",
                   "Defyne 0.5ml × 2 pre-jowl sulcus",
                   "המלצה: Touch-up אחרי שבועיים",
               ],
               outcome=[
                   "קו לסת חד וברור",
                   "סנטר מקבל הקרנה גברית",
                   "Pre-jowl נסגר ללא חוסר תנועה",
                   "Longevity צפוי 12–15 חודשים",
               ],
               color=NASHA_BLUE); page += 1

    slide_case(prs, page, total, 3,
               "מילוי שפתיים ראשוני",
               "אישה, גיל 28, ללא ניסיון קודם בפילר, מבקשת מילוי שפתיים 'טבעי'.",
               findings=[
                   "שפתיים דקות פיזיולוגית",
                   "אובדן vermillion border בעדינות",
                   "א-סימטריה מינורית בקשת קופידון",
                   "ציפיות 'אסתטיות, לא מוגזמות'",
               ],
               plan=[
                   "Kysse 0.5ml בלבד בטיפול ראשון",
                   "טכניקה: linear retrograde לאורך vermillion",
                   "Tenting נקודתי בקשת קופידון",
                   "מסר: 'פחות זה יותר' — touch-up בעוד שבועיים",
               ],
               outcome=[
                   "שיפור 'מורגש אך לא נראה'",
                   "תנועה טבעית, ללא 'בית מהפך' (lip flip מינורי)",
                   "מטופלת מבקשת touch-up של 0.3ml נוסף",
                   "תוצאה סופית: יציבה לכ-12 חודשים",
               ],
               color=OBT_PURPLE); page += 1

    # Closing
    slide_decision_tree(prs, page, total); page += 1
    slide_summary(prs, page, total); page += 1
    slide_thankyou(prs, page, total); page += 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Saved deck: {output_path}  ({page - 1} slides)")
    return page - 1, total


if __name__ == "__main__":
    out = Path(__file__).resolve().parent.parent / "Restylane_Clinical_Reference_AAA.pptx"
    built, expected = build(out)
    if built != expected:
        print(f"WARN: expected {expected} slides, built {built}")
